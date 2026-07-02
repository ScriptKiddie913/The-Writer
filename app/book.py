
from __future__ import annotations

import os
import io
import re
import sys
import json
import time
import uuid
import base64
import sqlite3
import logging
import threading
import traceback
from datetime import datetime
from contextlib import contextmanager
from typing import Optional, List, Dict, Any, Tuple

import requests
import numpy as np
from fastapi import FastAPI, UploadFile, File, Form
from fastapi.responses import HTMLResponse, JSONResponse, FileResponse, PlainTextResponse
from pydantic import BaseModel, Field, field_validator
import uvicorn

from reportlab.lib.pagesizes import LETTER
from reportlab.lib.units import inch
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.enums import TA_JUSTIFY, TA_CENTER, TA_LEFT
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.platypus import (
    Paragraph, Spacer, PageBreak, Frame, PageTemplate, BaseDocTemplate,
    Table, TableStyle, ListFlowable, ListItem,
)
from reportlab.lib import colors as rl_colors
from reportlab.pdfgen import canvas as pdfcanvas
from reportlab.graphics.shapes import Drawing, Rect, String, Line, Circle
from reportlab.graphics.charts.barcharts import VerticalBarChart
from reportlab.graphics.charts.lineplots import LinePlot, ScatterPlot
from reportlab.graphics.charts.piecharts import Pie

from document_profiles import get_document_profile, detect_document_type, MODE_TO_DEFAULT_DOC_TYPE, REPORT_DOC_TYPES

try:
    import fitz  # PyMuPDF — used to rasterize PDF pages for OCR
    HAS_FITZ = True
except Exception:
    HAS_FITZ = False

try:
    from pypdf import PdfReader
    HAS_PYPDF = True
except Exception:
    HAS_PYPDF = False

try:
    import docx as python_docx
    HAS_DOCX = True
except Exception:
    HAS_DOCX = False

try:
    from pptx import Presentation as python_pptx
    HAS_PPTX = True
except Exception:
    HAS_PPTX = False

try:
    import openpyxl
    HAS_OPENPYXL = True
except Exception:
    HAS_OPENPYXL = False

# =============================================================================
# CONFIG
# =============================================================================

OLLAMA_URL = os.environ.get("OLLAMA_URL", "http://localhost:11434")

def _get_base_dir() -> str:
    """
    Resolve the directory the app's persistent data should live in.

    - Normal `python book.py` run: folder containing this script.
    - Frozen PyInstaller --onefile exe: PyInstaller unpacks bundled code to a
      temp folder (sys._MEIPASS) that is wiped after the process exits, so we
      must NOT put the database/PDFs/logs there. Instead we use the folder
      that contains the actual .exe on disk, so data survives restarts.
    """
    if getattr(sys, "frozen", False):
        return os.path.dirname(os.path.abspath(sys.executable))
    return os.path.dirname(os.path.abspath(__file__))

BASE_DIR = _get_base_dir()
DB_PATH = os.environ.get("PRESS_DB_PATH", os.path.join(BASE_DIR, "the_writer.db"))
PDF_DIR = os.environ.get("PRESS_PDF_DIR", os.path.join(BASE_DIR, "generated_pdfs"))
LOG_PATH = os.environ.get("PRESS_LOG_PATH", os.path.join(BASE_DIR, "the_writer.log"))
PORT = int(os.environ.get("PRESS_PORT", "8000"))
DEFAULT_OCR_MODEL = os.environ.get("PRESS_OCR_MODEL", "glm-ocr:q8_0")
FONT_DIR = os.environ.get("PRESS_FONT_DIR", os.path.join(BASE_DIR, "fonts"))

os.makedirs(PDF_DIR, exist_ok=True)

MAX_UNITS = 80
MIN_UNITS = 1
MAX_PASSES_PER_UNIT = 14
OUTLINE_TOKEN_TARGET = 4500
HTTP_TIMEOUT = 900                  # local models on long context can be slow
OCR_TIMEOUT = 300
HEALTH_TIMEOUT = 8
MAX_OLLAMA_RETRIES = 3
RETRY_BACKOFF_SECONDS = 4
MAX_OUTLINE_PARSE_ATTEMPTS = 3
CONTINUITY_REWRITE_TRIGGER_WORDS = ("contradict", "inconsist", "error", "mismatch",
                                    "wrong", "conflict", "broken", "incorrect",
                                    "unsupported", "inaccurate", "off-tone", "off-brand")
EMBED_DIM_FALLBACK = 256
KB_CHUNK_WORDS = 220
KB_CHUNK_OVERLAP_WORDS = 30
KB_STRUCTURED_ROWS_PER_CHUNK = 20
PDF_RENDER_DPI = 150

# --- Writing modes ------------------------------------------------------
# unit         = the internal DB row (a "chapters" table row) generalized to
#                mean "one chapter / one story / one poem / one report section"
MODES = {
    "novel": {
        "label": "Novel",
        "unit": "Chapter", "unit_plural": "Chapters",
        "default_units": 18, "default_words_per_unit": 2600,
        "uses_story_bible": True, "uses_tables": False, "literary": True,
    },
    "short_story": {
        "label": "Short Story Collection",
        "unit": "Story", "unit_plural": "Stories",
        "default_units": 8, "default_words_per_unit": 1800,
        "uses_story_bible": True, "uses_tables": False, "literary": True,
    },
    "poetry": {
        "label": "Poetry Collection",
        "unit": "Poem", "unit_plural": "Poems",
        "default_units": 16, "default_words_per_unit": 220,
        "uses_story_bible": False, "uses_tables": False, "literary": True,
    },
    "report": {
        "label": "Executive Report",
        "unit": "Section", "unit_plural": "Sections",
        "default_units": 7, "default_words_per_unit": 650,
        "uses_story_bible": False, "uses_tables": True, "literary": False,
    },
}

# --- Visual themes --------------------------------------------------------
# Drive both PDF typography/color and a matching frontend accent color.
# Table/banner decor is only ever actually used when the project's mode
# also allows tables (report mode) — literary modes ignore that part.
THEMES = {
    "classic_cream": {
        "label": "Classic Cream", "css_accent": "#a4502b", "css_bg": "#f1e9d8",
        "css_bg2": "#e8dec7", "css_ink": "#2a241c",
        "pdf_body_font": "Times-Roman", "pdf_head_font": "Times-Bold",
        "pdf_accent": rl_colors.HexColor("#a4502b"), "pdf_ink": rl_colors.HexColor("#1a1a1a"),
        "pdf_muted": rl_colors.HexColor("#7a7266"), "pdf_page_bg": rl_colors.HexColor("#fffdf6"),
    },
    "midnight_gothic": {
        "label": "Midnight Gothic", "css_accent": "#8a5cd6", "css_bg": "#161421",
        "css_bg2": "#1f1c2e", "css_ink": "#e9e4f5",
        "pdf_body_font": "Times-Roman", "pdf_head_font": "Times-Bold",
        "pdf_accent": rl_colors.HexColor("#5c3a99"), "pdf_ink": rl_colors.HexColor("#151024"),
        "pdf_muted": rl_colors.HexColor("#6c6480"), "pdf_page_bg": rl_colors.HexColor("#ffffff"),
    },
    "rose_poetics": {
        "label": "Rose Poetics", "css_accent": "#b3446c", "css_bg": "#fbeff2",
        "css_bg2": "#f5dde4", "css_ink": "#3a1f27",
        "pdf_body_font": "Times-Italic", "pdf_head_font": "Times-Bold",
        "pdf_accent": rl_colors.HexColor("#b3446c"), "pdf_ink": rl_colors.HexColor("#2c161c"),
        "pdf_muted": rl_colors.HexColor("#8a6a72"), "pdf_page_bg": rl_colors.HexColor("#ffffff"),
    },
    "modern_corporate": {
        "label": "Modern Corporate", "css_accent": "#1f3a5f", "css_bg": "#eef1f5",
        "css_bg2": "#dde3ec", "css_ink": "#1b2431",
        "pdf_body_font": "Helvetica", "pdf_head_font": "Helvetica-Bold",
        "pdf_accent": rl_colors.HexColor("#1f3a5f"), "pdf_ink": rl_colors.HexColor("#1b2431"),
        "pdf_muted": rl_colors.HexColor("#5c6a7d"), "pdf_page_bg": rl_colors.HexColor("#ffffff"),
        "pdf_table_header_bg": rl_colors.HexColor("#1f3a5f"),
        "pdf_table_header_fg": rl_colors.white,
        "pdf_table_row_alt": rl_colors.HexColor("#f2f5f9"),
    },
    "minimal_mono": {
        "label": "Minimal Monochrome", "css_accent": "#333333", "css_bg": "#f5f5f3",
        "css_bg2": "#e6e6e2", "css_ink": "#1c1c1c",
        "pdf_body_font": "Helvetica", "pdf_head_font": "Helvetica-Bold",
        "pdf_accent": rl_colors.HexColor("#333333"), "pdf_ink": rl_colors.HexColor("#1c1c1c"),
        "pdf_muted": rl_colors.HexColor("#777777"), "pdf_page_bg": rl_colors.HexColor("#ffffff"),
        "pdf_table_header_bg": rl_colors.HexColor("#333333"),
        "pdf_table_header_fg": rl_colors.white,
        "pdf_table_row_alt": rl_colors.HexColor("#f0f0f0"),
    },
    "tactical_dark": {
        "label": "Tactical Dark", "css_accent": "#00ff9d", "css_bg": "#0a0e0f", "css_bg2": "#111518", "css_ink": "#c8d6d8",
        "pdf_body_font": "Helvetica", "pdf_head_font": "Helvetica-Bold",
        "pdf_accent": rl_colors.HexColor("#00c97a"), "pdf_ink": rl_colors.HexColor("#1a1a1a"),
        "pdf_muted": rl_colors.HexColor("#5a6a6e"), "pdf_page_bg": rl_colors.HexColor("#ffffff"),
        "pdf_table_header_bg": rl_colors.HexColor("#0d2a1f"),
        "pdf_table_header_fg": rl_colors.HexColor("#00c97a"),
        "pdf_table_row_alt": rl_colors.HexColor("#f0f7f4"),
    },
    "deep_navy": {
        "label": "Deep Navy", "css_accent": "#4fa3e0", "css_bg": "#0b1929", "css_bg2": "#112238", "css_ink": "#d8e8f5",
        "pdf_body_font": "Helvetica", "pdf_head_font": "Helvetica-Bold",
        "pdf_accent": rl_colors.HexColor("#1a5fa0"), "pdf_ink": rl_colors.HexColor("#111827"),
        "pdf_muted": rl_colors.HexColor("#4a6a8a"), "pdf_page_bg": rl_colors.HexColor("#ffffff"),
        "pdf_table_header_bg": rl_colors.HexColor("#1a3a5f"),
        "pdf_table_header_fg": rl_colors.white,
        "pdf_table_row_alt": rl_colors.HexColor("#eef4fb"),
    },
    "warm_editorial": {
        "label": "Warm Editorial", "css_accent": "#c0392b", "css_bg": "#faf6f0", "css_bg2": "#f2ebe0", "css_ink": "#2c1f14",
        "pdf_body_font": "Times-Roman", "pdf_head_font": "Times-Bold",
        "pdf_accent": rl_colors.HexColor("#c0392b"), "pdf_ink": rl_colors.HexColor("#1c130b"),
        "pdf_muted": rl_colors.HexColor("#8a6a5a"), "pdf_page_bg": rl_colors.HexColor("#fffdf8"),
    },
    "charcoal_ink": {
        "label": "Charcoal Ink", "css_accent": "#e8e0d0", "css_bg": "#1a1a1a", "css_bg2": "#242424", "css_ink": "#e0dbd0",
        "pdf_body_font": "Helvetica", "pdf_head_font": "Helvetica-Bold",
        "pdf_accent": rl_colors.HexColor("#555555"), "pdf_ink": rl_colors.HexColor("#111111"),
        "pdf_muted": rl_colors.HexColor("#7a7a7a"), "pdf_page_bg": rl_colors.HexColor("#ffffff"),
        "pdf_table_header_bg": rl_colors.HexColor("#222222"),
        "pdf_table_header_fg": rl_colors.white,
        "pdf_table_row_alt": rl_colors.HexColor("#f4f4f2"),
    },
}
DEFAULT_THEME_BY_MODE = {
    "novel": "classic_cream", "short_story": "classic_cream",
    "poetry": "rose_poetics", "report": "modern_corporate",
}

DOCUMENT_TYPE_CHOICES = [
    "novel",
    "short_story_collection",
    "poetry_collection",
    "executive_report",
    "cyber_threat_intelligence_report",
    "technical_manual",
    "academic_paper",
    "business_proposal",
    "market_research_report",
]

FONT_FALLBACKS = {
    "cormorant garamond": "Times-Bold",
    "eb garamond": "Times-Roman",
    "libre baskerville": "Times-Roman",
    "playfair display": "Times-Bold",
    "crimson text": "Times-Roman",
    "cinzel": "Times-Bold",
    "inter": "Helvetica",
    "manrope": "Helvetica",
    "source sans 3": "Helvetica",
    "rajdhani": "Helvetica-Bold",
    "ibm plex sans condensed": "Helvetica",
    "ibm plex sans": "Helvetica",
    "libertinus serif": "Times-Roman",
    "jetbrains mono": "Courier",
    "fira code": "Courier",
    "noto serif": "Times-Roman",
    "lora": "Times-Roman",
    "merriweather": "Times-Roman",
    "pt serif": "Times-Roman",
    "dm serif display": "Times-Bold",
    "space grotesk": "Helvetica",
    "dm sans": "Helvetica",
    "outfit": "Helvetica",
    "plus jakarta sans": "Helvetica",
    "geist mono": "Courier",
    "source code pro": "Courier",
    "inconsolata": "Courier",
}

CUSTOM_FONT_REGISTRY: Dict[str, str] = {}


def _font_key(name: str) -> str:
    return re.sub(r"[^a-z0-9]+", "", (name or "").lower())


def register_custom_fonts() -> None:
    if CUSTOM_FONT_REGISTRY:
        return
    if not os.path.isdir(FONT_DIR):
        return
    for entry in os.listdir(FONT_DIR):
        if not entry.lower().endswith((".ttf", ".otf")):
            continue
        font_name = os.path.splitext(entry)[0]
        font_path = os.path.join(FONT_DIR, entry)
        try:
            pdfmetrics.registerFont(TTFont(font_name, font_path))
            CUSTOM_FONT_REGISTRY[_font_key(font_name)] = font_name
        except Exception as exc:
            logger.warning("Could not register font '%s': %s", font_path, exc)


def resolve_pdf_font(preferred: str, role: str = "body") -> str:
    register_custom_fonts()
    preferred_key = _font_key(preferred)
    if preferred_key in CUSTOM_FONT_REGISTRY:
        return CUSTOM_FONT_REGISTRY[preferred_key]
    for key, registered in CUSTOM_FONT_REGISTRY.items():
        if preferred_key and preferred_key in key:
            return registered
    fallback_key = (preferred or "").strip().lower()
    if fallback_key in FONT_FALLBACKS:
        return FONT_FALLBACKS[fallback_key]
    if role in ("title", "heading"):
        return "Helvetica-Bold"
    if role == "mono":
        return "Courier"
    return "Helvetica"


def project_document_profile(proj: dict) -> Dict[str, Any]:
    return get_document_profile(
        proj.get("mode", "novel"),
        proj.get("document_type"),
        title=proj.get("title", ""),
        genre=proj.get("genre", ""),
        premise=proj.get("premise", ""),
        writing_notes=proj.get("writing_notes", ""),
    )

# =============================================================================
# LOGGING
# =============================================================================

logger = logging.getLogger("the_writer")
logger.setLevel(logging.INFO)
_fh = logging.FileHandler(LOG_PATH, encoding="utf-8")
_fh.setFormatter(logging.Formatter("%(asctime)s [%(levelname)s] %(message)s"))
_ch = logging.StreamHandler()
_ch.setFormatter(logging.Formatter("%(asctime)s [%(levelname)s] %(message)s"))
logger.addHandler(_fh)
logger.addHandler(_ch)
logger.propagate = False

# =============================================================================
# IN-MEMORY RUNTIME STATE
# =============================================================================

PROGRESS: Dict[str, Dict[str, Any]] = {}
STOP_FLAGS: Dict[str, bool] = {}
GEN_THREADS: Dict[str, threading.Thread] = {}
PROJECT_LOCKS: Dict[str, threading.Lock] = {}
PROGRESS_LOCK = threading.Lock()
KB_JOBS: Dict[str, Dict[str, Any]] = {}   # doc_id -> {"status":..., "error":...}


def project_lock(project_id: str) -> threading.Lock:
    with PROGRESS_LOCK:
        if project_id not in PROJECT_LOCKS:
            PROJECT_LOCKS[project_id] = threading.Lock()
        return PROJECT_LOCKS[project_id]


# =============================================================================
# DATABASE
# =============================================================================

def ensure_column_exists(conn: sqlite3.Connection, table: str, column: str, col_type: str) -> None:
    cursor = conn.execute(f"PRAGMA table_info({table})")
    columns = [row[1] for row in cursor.fetchall()]
    if column not in columns:
        conn.execute(f"ALTER TABLE {table} ADD COLUMN {column} {col_type}")
        logger.info(f"Added missing column '{column}' to table '{table}'")


def init_db() -> None:
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute("""
    CREATE TABLE IF NOT EXISTS projects (
        id TEXT PRIMARY KEY,
        title TEXT, mode TEXT DEFAULT 'novel', genre TEXT, premise TEXT,
        document_type TEXT DEFAULT '',
        theme TEXT DEFAULT 'classic_cream',
        writing_notes TEXT DEFAULT '',
        num_units INTEGER DEFAULT 10, words_per_unit INTEGER DEFAULT 2000,
        target_words INTEGER DEFAULT 20000,
        embed_model TEXT, ocr_model TEXT DEFAULT 'glm-ocr:q8_0',
        status TEXT DEFAULT 'draft',
        phase TEXT DEFAULT 'idle',
        total_units INTEGER DEFAULT 0,
        units_done INTEGER DEFAULT 0,
        words_written INTEGER DEFAULT 0,
        eta_seconds INTEGER DEFAULT 0,
        error TEXT,
        created_at TEXT, started_at TEXT, completed_at TEXT
    )""")
    c.execute("""
    CREATE TABLE IF NOT EXISTS agents (
        id TEXT PRIMARY KEY,
        project_id TEXT, name TEXT, model TEXT, role TEXT, order_idx INTEGER
    )""")
    c.execute("""
    CREATE TABLE IF NOT EXISTS chapters (
        id TEXT PRIMARY KEY,
        project_id TEXT, idx INTEGER, title TEXT, synopsis TEXT,
        key_events TEXT DEFAULT '[]',
        meta TEXT DEFAULT '{}',
        content TEXT DEFAULT '', summary TEXT DEFAULT '',
        word_count INTEGER DEFAULT 0, status TEXT DEFAULT 'pending',
        agent_name TEXT, revision_count INTEGER DEFAULT 0,
        continuity_notes TEXT DEFAULT ''
    )""")
    c.execute("""
    CREATE TABLE IF NOT EXISTS vectors (
        id TEXT PRIMARY KEY,
        project_id TEXT, chapter_idx INTEGER, kind TEXT,
        text TEXT, embedding BLOB, dim INTEGER, doc_id TEXT
    )""")
    c.execute("""
    CREATE TABLE IF NOT EXISTS logs (
        id TEXT PRIMARY KEY,
        project_id TEXT, ts TEXT, agent TEXT, role TEXT, message TEXT
    )""")
    c.execute("""
    CREATE TABLE IF NOT EXISTS entities (
        id TEXT PRIMARY KEY,
        project_id TEXT, name TEXT, name_key TEXT, type TEXT,
        facts TEXT DEFAULT '', first_chapter INTEGER, last_chapter INTEGER
    )""")
    c.execute("""
    CREATE TABLE IF NOT EXISTS knowledge_docs (
        id TEXT PRIMARY KEY,
        project_id TEXT, filename TEXT, source_type TEXT,
        char_count INTEGER DEFAULT 0, chunk_count INTEGER DEFAULT 0,
        status TEXT DEFAULT 'processing', error TEXT DEFAULT '',
        created_at TEXT
    )""")
    c.execute("CREATE INDEX IF NOT EXISTS idx_chapters_proj ON chapters(project_id, idx)")
    c.execute("CREATE INDEX IF NOT EXISTS idx_vectors_proj ON vectors(project_id)")
    c.execute("CREATE INDEX IF NOT EXISTS idx_vectors_doc ON vectors(doc_id)")
    c.execute("CREATE INDEX IF NOT EXISTS idx_logs_proj ON logs(project_id, ts)")
    c.execute("CREATE INDEX IF NOT EXISTS idx_entities_proj ON entities(project_id, name_key)")
    c.execute("CREATE INDEX IF NOT EXISTS idx_kb_proj ON knowledge_docs(project_id)")

    # --- Migration: ensure all columns that may be missing from older schemas ---
    ensure_column_exists(conn, "chapters", "key_events", "TEXT DEFAULT '[]'")
    ensure_column_exists(conn, "chapters", "meta", "TEXT DEFAULT '{}'")
    ensure_column_exists(conn, "chapters", "revision_count", "INTEGER DEFAULT 0")
    ensure_column_exists(conn, "chapters", "continuity_notes", "TEXT DEFAULT ''")
    ensure_column_exists(conn, "chapters", "synopsis", "TEXT")
    ensure_column_exists(conn, "chapters", "content", "TEXT DEFAULT ''")
    ensure_column_exists(conn, "chapters", "summary", "TEXT DEFAULT ''")
    ensure_column_exists(conn, "chapters", "word_count", "INTEGER DEFAULT 0")
    ensure_column_exists(conn, "chapters", "status", "TEXT DEFAULT 'pending'")
    ensure_column_exists(conn, "chapters", "agent_name", "TEXT")
    ensure_column_exists(conn, "entities", "facts", "TEXT DEFAULT ''")
    ensure_column_exists(conn, "entities", "first_chapter", "INTEGER")
    ensure_column_exists(conn, "entities", "last_chapter", "INTEGER")
    ensure_column_exists(conn, "vectors", "doc_id", "TEXT")
    ensure_column_exists(conn, "projects", "mode", "TEXT DEFAULT 'novel'")
    ensure_column_exists(conn, "projects", "document_type", "TEXT DEFAULT ''")
    ensure_column_exists(conn, "projects", "theme", "TEXT DEFAULT 'classic_cream'")
    ensure_column_exists(conn, "projects", "writing_notes", "TEXT DEFAULT ''")
    ensure_column_exists(conn, "projects", "num_units", "INTEGER DEFAULT 10")
    ensure_column_exists(conn, "projects", "words_per_unit", "INTEGER DEFAULT 2000")
    ensure_column_exists(conn, "projects", "ocr_model", "TEXT DEFAULT 'glm-ocr:q8_0'")
    ensure_column_exists(conn, "projects", "total_units", "INTEGER DEFAULT 0")
    ensure_column_exists(conn, "projects", "units_done", "INTEGER DEFAULT 0")

    conn.commit()
    conn.close()
    logger.info("Database initialized at %s", DB_PATH)


@contextmanager
def db():
    conn = sqlite3.connect(DB_PATH, timeout=30)
    conn.row_factory = sqlite3.Row
    try:
        yield conn
        conn.commit()
    except Exception:
        conn.rollback()
        raise
    finally:
        conn.close()


def get_project(project_id: str) -> Optional[Dict[str, Any]]:
    with db() as conn:
        row = conn.execute("SELECT * FROM projects WHERE id=?", (project_id,)).fetchone()
        return dict(row) if row else None


def list_projects() -> List[Dict[str, Any]]:
    with db() as conn:
        rows = conn.execute(
            "SELECT id, title, mode, document_type, genre, status, phase, theme, num_units, "
            "units_done, total_units, words_written, target_words, created_at "
            "FROM projects ORDER BY created_at DESC"
        ).fetchall()
    return [dict(r) for r in rows]


def delete_project(project_id: str) -> None:
    with db() as conn:
        for table in ("projects", "agents", "chapters", "vectors", "logs",
                       "entities", "knowledge_docs"):
            conn.execute(f"DELETE FROM {table} WHERE {'id' if table=='projects' else 'project_id'}=?",
                         (project_id,))
    pdf_path = os.path.join(PDF_DIR, f"{project_id}.pdf")
    if os.path.exists(pdf_path):
        os.remove(pdf_path)
    PROGRESS.pop(project_id, None)
    STOP_FLAGS.pop(project_id, None)


def get_agents(project_id: str) -> List[Dict[str, Any]]:
    with db() as conn:
        rows = conn.execute(
            "SELECT * FROM agents WHERE project_id=? ORDER BY order_idx", (project_id,)
        ).fetchall()
        return [dict(r) for r in rows]


def get_chapters(project_id: str) -> List[Dict[str, Any]]:
    with db() as conn:
        rows = conn.execute(
            "SELECT * FROM chapters WHERE project_id=? ORDER BY idx", (project_id,)
        ).fetchall()
        return [dict(r) for r in rows]


def get_chapter(project_id: str, idx: int) -> Optional[Dict[str, Any]]:
    with db() as conn:
        row = conn.execute(
            "SELECT * FROM chapters WHERE project_id=? AND idx=?", (project_id, idx)
        ).fetchone()
        return dict(row) if row else None


def save_chapter_content(project_id: str, idx: int, content: str) -> None:
    with db() as conn:
        conn.execute(
            "UPDATE chapters SET content=?, word_count=? WHERE project_id=? AND idx=?",
            (content, word_count(content), project_id, idx)
        )


def total_words_written(project_id: str) -> int:
    with db() as conn:
        row = conn.execute(
            "SELECT COALESCE(SUM(word_count),0) AS w FROM chapters WHERE project_id=?",
            (project_id,)
        ).fetchone()
        return row["w"]


# =============================================================================
# OLLAMA CLIENT  (retries + health checks; no external APIs ever touched)
# =============================================================================

class OllamaError(Exception):
    pass


def ollama_health() -> Tuple[bool, str]:
    try:
        r = requests.get(f"{OLLAMA_URL}/api/tags", timeout=HEALTH_TIMEOUT)
        r.raise_for_status()
        return True, "reachable"
    except requests.exceptions.ConnectionError:
        return False, f"Cannot connect to Ollama at {OLLAMA_URL}. Is 'ollama serve' running?"
    except Exception as e:
        return False, f"Ollama health check failed: {e}"


def ollama_list_models() -> List[str]:
    try:
        r = requests.get(f"{OLLAMA_URL}/api/tags", timeout=HEALTH_TIMEOUT)
        r.raise_for_status()
        data = r.json()
        return sorted([m["name"] for m in data.get("models", [])])
    except Exception as e:
        logger.warning("Could not list Ollama models: %s", e)
        return []


def ollama_chat(model: str, system: str, user: str, temperature: float = 0.92,
                 num_predict: int = 950, repeat_penalty: float = 1.15,
                 num_ctx: int = 8192, retries: int = MAX_OLLAMA_RETRIES) -> str:
    """Call the local Ollama chat endpoint with retry/backoff. Raises OllamaError
    only after all retries are exhausted, so callers can surface a clean message
    instead of a stack trace mid-book."""
    payload = {
        "model": model,
        "messages": [
            {"role": "system", "content": system},
            {"role": "user", "content": user},
        ],
        "stream": False,
        "options": {
            "temperature": temperature,
            "num_predict": num_predict,
            "repeat_penalty": repeat_penalty,
            "num_ctx": num_ctx,
        },
    }
    last_err = None
    for attempt in range(1, retries + 1):
        try:
            r = requests.post(f"{OLLAMA_URL}/api/chat", json=payload, timeout=HTTP_TIMEOUT)
            r.raise_for_status()
            data = r.json()
            content = data.get("message", {}).get("content", "").strip()
            if not content:
                raise OllamaError(f"Model '{model}' returned empty content.")
            return content
        except requests.exceptions.ConnectionError as e:
            last_err = f"Connection error talking to Ollama ({model}): {e}"
        except requests.exceptions.Timeout as e:
            last_err = f"Timeout waiting on Ollama ({model}) after {HTTP_TIMEOUT}s: {e}"
        except requests.exceptions.HTTPError as e:
            body = ""
            try:
                body = r.text[:300]
            except Exception:
                pass
            last_err = f"Ollama HTTP error for model '{model}': {e} {body}"
        except Exception as e:
            last_err = f"Unexpected error calling Ollama model '{model}': {e}"

        logger.warning("ollama_chat attempt %d/%d failed: %s", attempt, retries, last_err)
        if attempt < retries:
            time.sleep(RETRY_BACKOFF_SECONDS * attempt)

    raise OllamaError(last_err or "ollama_chat failed for unknown reasons.")


def ollama_embed(model: str, text: str) -> np.ndarray:
    try:
        r = requests.post(f"{OLLAMA_URL}/api/embeddings",
                           json={"model": model, "prompt": text[:6000]}, timeout=60)
        r.raise_for_status()
        emb = r.json().get("embedding")
        if emb:
            return np.array(emb, dtype=np.float32)
    except Exception as e:
        logger.debug("Embedding call failed (%s) — falling back to local hashing embed.", e)
    return hash_embed(text)


def hash_embed(text: str, dim: int = EMBED_DIM_FALLBACK) -> np.ndarray:
    """Deterministic, fully local fallback embedding. Guarantees RAG keeps working
    even if no embedding model is installed in Ollama — no external API is ever
    used for this."""
    vec = np.zeros(dim, dtype=np.float32)
    for tok in re.findall(r"[a-zA-Z']+", text.lower()):
        h = hash(tok) % dim
        vec[h] += 1.0
    norm = np.linalg.norm(vec)
    return vec / norm if norm > 0 else vec


def ollama_ocr_image(model: str, image_b64: str, prompt: str,
                      retries: int = MAX_OLLAMA_RETRIES) -> str:
    """Runs a local vision/OCR model (e.g. glm-ocr:q8_0) over a single base64
    PNG/JPEG image via Ollama's /api/generate endpoint. Never touches any
    external service — OCR happens entirely against the local Ollama daemon."""
    payload = {
        "model": model,
        "prompt": prompt,
        "images": [image_b64],
        "stream": False,
        "options": {"temperature": 0.1, "num_predict": 4000},
    }
    last_err = None
    for attempt in range(1, retries + 1):
        try:
            r = requests.post(f"{OLLAMA_URL}/api/generate", json=payload, timeout=OCR_TIMEOUT)
            r.raise_for_status()
            data = r.json()
            text = (data.get("response") or "").strip()
            return text
        except requests.exceptions.ConnectionError as e:
            last_err = f"Connection error talking to Ollama OCR ({model}): {e}"
        except requests.exceptions.Timeout as e:
            last_err = f"Timeout waiting on Ollama OCR ({model}) after {OCR_TIMEOUT}s: {e}"
        except Exception as e:
            last_err = f"Unexpected error calling Ollama OCR model '{model}': {e}"
        logger.warning("ollama_ocr_image attempt %d/%d failed: %s", attempt, retries, last_err)
        if attempt < retries:
            time.sleep(RETRY_BACKOFF_SECONDS * attempt)
    raise OllamaError(last_err or "ollama_ocr_image failed for unknown reasons.")


# =============================================================================
# VECTOR STORE / RAG   (project-scoped — a book's vectors are only ever
# retrieved by queries carrying that same project_id, so per-book knowledge
# bases and story memories can never leak into another book's generation.)
# =============================================================================

def vector_add(project_id: str, chapter_idx: int, kind: str, text: str,
                embed_model: str, doc_id: Optional[str] = None) -> None:
    emb = ollama_embed(embed_model, text)
    with db() as conn:
        conn.execute(
            "INSERT INTO vectors (id, project_id, chapter_idx, kind, text, embedding, dim, doc_id) "
            "VALUES (?,?,?,?,?,?,?,?)",
            (str(uuid.uuid4()), project_id, chapter_idx, kind, text, emb.tobytes(),
             emb.shape[0], doc_id)
        )


def vector_retrieve(project_id: str, query_text: str, embed_model: str,
                     top_k: int = 6, kind: Optional[str] = None) -> List[Tuple[float, int, str, str]]:
    with db() as conn:
        if kind:
            rows = conn.execute(
                "SELECT chapter_idx, kind, text, embedding, dim FROM vectors "
                "WHERE project_id=? AND kind=?", (project_id, kind)
            ).fetchall()
        else:
            rows = conn.execute(
                "SELECT chapter_idx, kind, text, embedding, dim FROM vectors WHERE project_id=?",
                (project_id,)
            ).fetchall()
    if not rows:
        return []
    q = ollama_embed(embed_model, query_text)
    qn = float(np.linalg.norm(q))
    scored = []
    for row in rows:
        v = np.frombuffer(row["embedding"], dtype=np.float32)
        if v.shape[0] != q.shape[0]:
            continue
        vn = float(np.linalg.norm(v))
        sim = float(np.dot(q, v) / (qn * vn)) if (qn > 0 and vn > 0) else 0.0
        scored.append((sim, row["chapter_idx"], row["kind"], row["text"]))
    scored.sort(key=lambda x: x[0], reverse=True)
    return scored[:top_k]


# =============================================================================
# KNOWLEDGE BASE — per-book PDFs / images / pasted text
# =============================================================================
# Strictly scoped: every knowledge_docs row and every vectors row it produces
# carries this book's project_id, and vector_retrieve() always filters on
# project_id, so a book's knowledge base cannot bleed into any other book's
# generation, ever.

def chunk_text(text: str, chunk_words: int = KB_CHUNK_WORDS,
               overlap_words: int = KB_CHUNK_OVERLAP_WORDS) -> List[str]:
    words = text.split()
    if not words:
        return []
    chunks = []
    step = max(chunk_words - overlap_words, 1)
    for start in range(0, len(words), step):
        chunk = " ".join(words[start:start + chunk_words])
        if chunk.strip():
            chunks.append(chunk.strip())
        if start + chunk_words >= len(words):
            break
    return chunks


def extract_text_from_pdf(pdf_bytes: bytes, ocr_model: str, project_id: str,
                           filename: str) -> str:
    """Rasterizes each PDF page and OCRs it with the local vision model if
    PyMuPDF is available; otherwise falls back to extracting the embedded
    text layer only (works for text-based PDFs, not scans)."""
    if HAS_FITZ:
        text_parts = []
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        try:
            for page_num in range(len(doc)):
                page = doc[page_num]
                pix = page.get_pixmap(dpi=PDF_RENDER_DPI)
                img_bytes = pix.tobytes("png")
                b64 = base64.b64encode(img_bytes).decode("ascii")
                prompt = ("Transcribe every word of visible text on this document page "
                          "exactly as written, preserving reading order and paragraph/line "
                          "breaks. Include any table content as plain rows. Output only the "
                          "transcribed text, no commentary.")
                try:
                    page_text = ollama_ocr_image(ocr_model, b64, prompt)
                except OllamaError as e:
                    log(project_id, "system", "knowledge",
                        f"OCR failed on page {page_num+1} of '{filename}': {e}")
                    page_text = ""
                if page_text:
                    text_parts.append(page_text)
        finally:
            doc.close()
        combined = "\n\n".join(text_parts).strip()
        if combined:
            return combined
        # fall through to text-layer extraction if OCR produced nothing
    if HAS_PYPDF:
        try:
            reader = PdfReader(io.BytesIO(pdf_bytes))
            parts = [(p.extract_text() or "") for p in reader.pages]
            combined = "\n\n".join(parts).strip()
            if combined:
                return combined
        except Exception as e:
            log(project_id, "system", "knowledge", f"pypdf text extraction failed for "
                                                     f"'{filename}': {e}")
    raise OllamaError(
        "Could not extract any text from this PDF. Install PyMuPDF for OCR of scanned "
        "PDFs ('pip install pymupdf'), or paste the text directly instead."
    )


def extract_text_from_image(image_bytes: bytes, ocr_model: str) -> str:
    b64 = base64.b64encode(image_bytes).decode("ascii")
    prompt = ("Transcribe every word of visible text in this image exactly as written, "
              "preserving reading order and line breaks. Output only the transcribed "
              "text, no commentary.")
    return ollama_ocr_image(ocr_model, b64, prompt)


def extract_text_from_csv(raw_bytes: bytes) -> str:
    """Parse CSV/TSV bytes into a readable plain-text table string."""
    import csv, io as _io
    text_io = _io.StringIO(raw_bytes.decode("utf-8", errors="ignore"))
    sample = text_io.read(2048)
    dialect = "excel-tab" if "\t" in sample else "excel"
    text_io.seek(0)
    reader = csv.reader(text_io, dialect=dialect)
    rows = list(reader)
    if not rows:
        return ""
    col_widths = [max(len(str(r[i])) for r in rows if i < len(r)) for i in range(len(rows[0]))]
    lines = []
    for row in rows:
        lines.append("  ".join(str(row[i]).ljust(col_widths[i]) for i in range(len(row))))
    return "\n".join(lines)


def extract_text_from_xlsx(raw_bytes: bytes) -> str:
    """Extract text from Excel workbooks using openpyxl (all sheets)."""
    if not HAS_OPENPYXL:
        raise OllamaError("openpyxl is not installed. Run: pip install openpyxl --break-system-packages")
    wb = openpyxl.load_workbook(io.BytesIO(raw_bytes), read_only=True, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"=== Sheet: {sheet_name} ===")
        for row in ws.iter_rows(values_only=True):
            row_text = "  ".join(str(c) if c is not None else "" for c in row).rstrip()
            if row_text.strip():
                parts.append(row_text)
    return "\n".join(parts)


def extract_text_from_docx(raw_bytes: bytes) -> str:
    """Extract text from .docx files using python-docx."""
    if not HAS_DOCX:
        raise OllamaError("python-docx is not installed. Run: pip install python-docx --break-system-packages")
    doc = python_docx.Document(io.BytesIO(raw_bytes))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            parts.append("  ".join(cell.text.strip() for cell in row.cells))
    return "\n".join(parts)


def extract_text_from_pptx(raw_bytes: bytes) -> str:
    """Extract all slide text from .pptx files using python-pptx."""
    if not HAS_PPTX:
        raise OllamaError("python-pptx is not installed. Run: pip install python-pptx --break-system-packages")
    prs = python_pptx(io.BytesIO(raw_bytes))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"=== Slide {i} ===")
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        parts.append(t)
        if hasattr(slide, "notes_slide") and slide.notes_slide:
            try:
                note_text = slide.notes_slide.notes_text_frame.text.strip()
                if note_text:
                    parts.append(f"[Notes] {note_text}")
            except Exception:
                pass
    return "\n".join(parts)


def extract_text_from_json(raw_bytes: bytes) -> str:
    """Pretty-print JSON as plain text so it can be chunked and embedded."""
    try:
        data = json.loads(raw_bytes.decode("utf-8", errors="ignore"))
        return json.dumps(data, indent=2, ensure_ascii=False)
    except Exception:
        return raw_bytes.decode("utf-8", errors="ignore")


def _chunk_by_lines(text: str, lines_per_chunk: int = KB_STRUCTURED_ROWS_PER_CHUNK) -> List[str]:
    lines = [ln for ln in (text or "").splitlines() if ln.strip()]
    chunks = []
    for start in range(0, len(lines), lines_per_chunk):
        chunk = "\n".join(lines[start:start + lines_per_chunk]).strip()
        if chunk:
            chunks.append(chunk)
    return chunks


def chunk_by_source_type(text: str, source_type: str) -> List[str]:
    source_type = (source_type or "").lower()
    if source_type in ("csv", "xlsx"):
        return _chunk_by_lines(text, lines_per_chunk=KB_STRUCTURED_ROWS_PER_CHUNK)
    if source_type == "pptx":
        return [s.strip() for s in text.split("=== Slide") if s.strip()]
    if source_type == "docx":
        return [s.strip() for s in re.split(r"\n\s*\n", text) if s.strip()]
    if source_type == "json":
        try:
            data = json.loads(text)
            if isinstance(data, list):
                return [json.dumps(item, ensure_ascii=False, indent=2) for item in data]
            if isinstance(data, dict):
                return [f"{k}: {json.dumps(v, ensure_ascii=False, indent=2)}" for k, v in data.items()]
        except Exception:
            pass
    return chunk_text(text)


def _source_tag(filename: str, idx: int, total: int) -> str:
    return f"[SOURCE: {filename} | CHUNK: {idx}/{total}]"


def _tagged_chunks(filename: str, chunks: List[str]) -> List[str]:
    total = max(len(chunks), 1)
    return [f"{_source_tag(filename, i + 1, total)}\n{chunk}" for i, chunk in enumerate(chunks)]


def ingest_knowledge_doc(project_id: str, doc_id: str, filename: str, source_type: str,
                          raw_bytes: Optional[bytes], raw_text: Optional[str],
                          embed_model: str, ocr_model: str) -> None:
    """Runs in a background thread. Extracts (if needed), chunks, embeds, and
    stores knowledge text for exactly one project. Never touches any other
    project's data."""
    try:
        if source_type == "pdf":
            text = extract_text_from_pdf(raw_bytes, ocr_model, project_id, filename)
        elif source_type == "image":
            text = extract_text_from_image(raw_bytes, ocr_model)
        elif source_type == "text_file":
            text = raw_bytes.decode("utf-8", errors="ignore")
        elif source_type == "markdown":
            text = raw_bytes.decode("utf-8", errors="ignore")
        elif source_type == "csv":
            text = extract_text_from_csv(raw_bytes)
        elif source_type == "xlsx":
            text = extract_text_from_xlsx(raw_bytes)
        elif source_type == "docx":
            text = extract_text_from_docx(raw_bytes)
        elif source_type == "pptx":
            text = extract_text_from_pptx(raw_bytes)
        elif source_type == "json":
            text = extract_text_from_json(raw_bytes)
        elif source_type == "paste":
            text = raw_text or ""
        else:
            raise OllamaError(f"Unknown knowledge source type: {source_type}")

        text = text.strip()
        if not text:
            raise OllamaError("No text could be extracted from this document.")

        chunks = chunk_by_source_type(text, source_type)
        tagged_chunks = _tagged_chunks(filename, chunks)
        for chunk in tagged_chunks:
            vector_add(project_id, -1, "knowledge", chunk, embed_model, doc_id=doc_id)

        with db() as conn:
            conn.execute(
                "UPDATE knowledge_docs SET status='ready', char_count=?, chunk_count=? WHERE id=?",
                (len(text), len(tagged_chunks), doc_id)
            )
        log(project_id, "system", "knowledge",
            f"Knowledge doc '{filename}' ready: {len(text)} chars, {len(tagged_chunks)} chunks.")
    except Exception as e:
        with db() as conn:
            conn.execute(
                "UPDATE knowledge_docs SET status='error', error=? WHERE id=?",
                (str(e)[:500], doc_id)
            )
        log(project_id, "system", "error", f"Knowledge ingestion failed for '{filename}': {e}")


def knowledge_context(project_id: str, query_text: str, embed_model: str,
                       top_k: int = 5) -> str:
    retrieved = vector_retrieve(project_id, query_text, embed_model, top_k=top_k, kind="knowledge")
    if not retrieved:
        return ""
    return "\n".join(f"- {r[3]}" for r in retrieved)


# =============================================================================
# STORY BIBLE — structured entity tracking (characters / locations / objects)
# Used only for novel and short_story modes.
# =============================================================================

ENTITY_EXTRACT_SYS = (
    "You are a story-bible extraction engine. You read a piece of fiction and output "
    "STRICT JSON describing every named character, location, and significant "
    "object/technology that appears or is updated in it. No commentary, no markdown fences."
)

def entity_extract_prompt(chapter: dict, chapter_text: str) -> str:
    return f"""Read this text and extract its story-bible entries.

TITLE: {chapter['title']}
TEXT:
{chapter_text[:7000]}

For every named character, location, or significant object/technology, output one entry.
Facts should be short, concrete, and cumulative-safe (state things that are now true).

Output STRICT JSON: a list of objects shaped exactly like:
{{"name": "...", "type": "character|location|object", "facts": "short factual notes"}}

Output ONLY the JSON array. If nothing notable is introduced or changed, output [].""".strip()


def _name_key(name: str) -> str:
    return re.sub(r"[^a-z0-9]+", "", name.lower())


def extract_and_merge_entities(project_id: str, chapter: dict, chapter_text: str,
                                extractor_model: str) -> int:
    try:
        raw = ollama_chat(extractor_model, ENTITY_EXTRACT_SYS,
                           entity_extract_prompt(chapter, chapter_text),
                           temperature=0.2, num_predict=1200)
        entries = safe_json_array(raw) if "[" in raw else []
    except Exception as e:
        log(project_id, "system", "error", f"Entity extraction skipped for "
                                             f"'{chapter['title']}': {e}")
        return 0

    touched = 0
    with db() as conn:
        for entry in entries:
            if not isinstance(entry, dict) or not entry.get("name"):
                continue
            name = str(entry["name"]).strip()[:120]
            etype = str(entry.get("type", "character")).strip().lower()
            if etype not in ("character", "location", "object"):
                etype = "character"
            facts = str(entry.get("facts", "")).strip()[:600]
            key = _name_key(name)
            if not key:
                continue
            existing = conn.execute(
                "SELECT id, facts FROM entities WHERE project_id=? AND name_key=?",
                (project_id, key)
            ).fetchone()
            if existing:
                merged_facts = (existing["facts"] + " | " + facts).strip(" |") if facts else existing["facts"]
                if len(merged_facts) > 1200:
                    merged_facts = merged_facts[-1200:]
                conn.execute(
                    "UPDATE entities SET facts=?, last_chapter=? WHERE id=?",
                    (merged_facts, chapter["idx"], existing["id"])
                )
            else:
                conn.execute(
                    "INSERT INTO entities (id, project_id, name, name_key, type, facts, "
                    "first_chapter, last_chapter) VALUES (?,?,?,?,?,?,?,?)",
                    (str(uuid.uuid4()), project_id, name, key, etype, facts,
                     chapter["idx"], chapter["idx"])
                )
            touched += 1
    return touched


def get_story_bible(project_id: str, limit: int = 60) -> List[Dict[str, Any]]:
    with db() as conn:
        rows = conn.execute(
            "SELECT name, type, facts, first_chapter, last_chapter FROM entities "
            "WHERE project_id=? ORDER BY last_chapter DESC, name ASC LIMIT ?",
            (project_id, limit)
        ).fetchall()
    return [dict(r) for r in rows]


def story_bible_context(project_id: str, max_entities: int = 25) -> str:
    bible = get_story_bible(project_id, limit=max_entities)
    if not bible:
        return ""
    lines = [f"- [{e['type'].upper()}] {e['name']}: {e['facts']}" for e in bible if e["facts"]]
    return "\n".join(lines)


# =============================================================================
# LOGGING / PROGRESS HELPERS
# =============================================================================

def log(project_id: str, agent: str, role: str, message: str) -> None:
    with db() as conn:
        conn.execute(
            "INSERT INTO logs (id, project_id, ts, agent, role, message) VALUES (?,?,?,?,?,?)",
            (str(uuid.uuid4()), project_id, datetime.utcnow().isoformat(), agent, role, message[:4000])
        )
    logger.info("[%s] %s/%s: %s", project_id[:8], agent, role, message[:200])


def update_progress(project_id: str, **kwargs) -> None:
    with PROGRESS_LOCK:
        p = PROGRESS.setdefault(project_id, {})
        p.update(kwargs)
    persistable = {k: v for k, v in kwargs.items() if k in (
        "status", "phase", "total_units", "units_done",
        "words_written", "eta_seconds", "error"
    )}
    if persistable:
        sets = ", ".join(f"{k}=?" for k in persistable)
        vals = list(persistable.values()) + [project_id]
        with db() as conn:
            conn.execute(f"UPDATE projects SET {sets} WHERE id=?", vals)


def word_count(text: str) -> int:
    return len(text.split()) if text else 0


# =============================================================================
# PROMPTS — mode-dispatched, functional roles only, no personas
# =============================================================================

def style_notes_block(proj: dict) -> str:
    notes = (proj.get("writing_notes") or "").strip()
    if not notes:
        return ""
    return f"\nAUTHOR / HOUSE STYLE NOTES (follow these closely):\n{notes}\n"


# ---- Architect ------------------------------------------------------------

def architect_sys_for_mode(mode: str, doc_profile: Optional[dict] = None) -> str:
    if mode == "report":
        doc_label = (doc_profile or {}).get("label", "report")
        return (
            f"You are a {doc_label.lower()}-structuring engine. You design the section-by-section "
            "outline for a professional document for an autonomous writing system. You output STRICT JSON ONLY. "
            "No markdown fences, no commentary, no preamble, no trailing text of any kind."
        )
    if mode == "poetry":
        return (
            "You are a poetry-collection curator engine. You design the poem-by-poem "
            "table of contents for an autonomous writing system. You output STRICT JSON "
            "ONLY. No markdown fences, no commentary, no preamble, no trailing text."
        )
    return (
        "You are a narrative-architecture engine. You design detailed, original outlines "
        "for an autonomous writing system. You output STRICT JSON ONLY. No markdown "
        "fences, no commentary, no preamble, no trailing text of any kind."
    )


def architect_prompt(mode: str, title: str, genre: str, premise: str, num_units: int,
                      words_per_unit: int, proj: dict, kb_context: str = "") -> str:
    m = MODES[mode]
    doc_profile = project_document_profile(proj)
    kb_block = f"\nRELEVANT KNOWLEDGE BASE MATERIAL PROVIDED BY THE USER:\n{kb_context}\n" if kb_context else ""
    style_block = style_notes_block(proj)

    if mode == "novel":
        return f"""Design a complete chapter-by-chapter outline for a full-length {genre} novel.

TITLE: {title}
PREMISE: {premise}
TOTAL CHAPTERS REQUIRED: {num_units}
TARGET LENGTH PER CHAPTER: ~{words_per_unit} words
{style_block}{kb_block}
Requirements:
- A coherent plot arc with setup, rising action, a midpoint turn, climax, and resolution
  spread across the {num_units} chapters.
- Concrete, specific events per chapter — not vague themes or moods.
- Consistent character names, locations, and world/technology rules established early
  and reused consistently throughout.
- This is a serious, mature, fully-developed novel — not a synopsis sketch.

Output STRICT JSON: a list of exactly {num_units} objects, each shaped exactly like:
{{"title": "...", "synopsis": "2-4 sentences of concrete plot content", "key_events": ["...", "...", "..."], "pov_character": "name of the POV character for this chapter", "setting": "primary location/time of this chapter", "tone": "e.g. tense, melancholic, comic, action-driven", "epigraph": "optional: a short quote or lines to open the chapter (leave empty string if none)"}}

Output ONLY the JSON array, nothing else.""".strip()

    if mode == "short_story":
        return f"""Design a table of contents for a collection of {num_units} standalone
short stories in the {genre} genre/tradition, unified by the premise below.

COLLECTION TITLE: {title}
PREMISE / UNIFYING IDEA: {premise}
TARGET LENGTH PER STORY: ~{words_per_unit} words
{style_block}{kb_block}
Requirements:
- Each story must be complete in itself (its own beginning, middle, end) while sharing
  the collection's tone, world, or theme as appropriate to the premise.
- Concrete, specific premises per story — not vague themes.
- Vary structure/angle story to story so the collection doesn't feel repetitive.

Output STRICT JSON: a list of exactly {num_units} objects, each shaped exactly like:
{{"title": "...", "synopsis": "2-4 sentences of concrete story premise", "key_events": ["...", "...", "..."]}}

Output ONLY the JSON array, nothing else.""".strip()

    if mode == "poetry":
        return f"""Design a table of contents for a poetry collection of {num_units} poems
in the {genre} tradition/style, unified by the premise below.

COLLECTION TITLE: {title}
PREMISE / UNIFYING THEME: {premise}
TARGET LENGTH PER POEM: ~{words_per_unit} words
{style_block}{kb_block}
Requirements:
- Vary form across the collection (e.g. free verse, sonnet, villanelle, haiku sequence,
  prose poem) where it fits the premise — do not make every poem the same shape.
- Each poem needs its own concrete subject/image/angle, not just a mood restated.

Output STRICT JSON: a list of exactly {num_units} objects, each shaped exactly like:
{{"title": "...", "synopsis": "concrete subject, imagery, and emotional arc of the poem", "form": "e.g. free verse / sonnet / haiku sequence / villanelle / prose poem", "key_events": []}}

Output ONLY the JSON array, nothing else.""".strip()

    content_rules = "\n".join(f"- {rule}" for rule in doc_profile.get("content_rules", []))
    allowed_blocks = ", ".join(doc_profile.get("allowed_blocks", []))
    return f"""Design the section-by-section structure of a professional {doc_profile['label'].lower()}.

REPORT TITLE: {title}
DOMAIN / INDUSTRY: {genre}
PURPOSE / BRIEF: {premise}
TOTAL SECTIONS REQUIRED: {num_units}
TARGET LENGTH PER SECTION: ~{words_per_unit} words
{style_block}{kb_block}
Requirements:
- Use a structure that fits this document type rather than forcing a generic report template.
- Hard guardrails for this document type:
{content_rules}
- Supported visual/content blocks for rendering: {allowed_blocks}
- Each section needs a concrete synopsis of exactly what it must cover and, where the
    knowledge base or premise implies comparable data (numbers, options, timelines,
    metrics), mark it as needing a visual/data block.

Output STRICT JSON: a list of exactly {num_units} objects, each shaped exactly like:
{{"title": "...", "synopsis": "concrete description of what this section must cover", "include_table": true or false, "table_spec": "if include_table, describe the table's columns and the data it should contain; else empty string", "include_chart": true or false, "chart_spec": "if include_chart, describe the chart type and its data; else empty string", "key_events": []}}

Output ONLY the JSON array, nothing else.""".strip()


# ---- Writer -----------------------------------------------------------------

def writer_sys_for_mode(mode: str, doc_profile: Optional[dict] = None) -> str:
    doc_profile = doc_profile or {}
    if mode == "poetry":
        return (
            "You are a professional poetry-writing engine operating with full creative "
            "latitude, constrained only by the collection context given to you. You write "
            "vivid, precise, well-crafted poems in the requested form, with attention to "
            "imagery, sound, line breaks, and emotional arc. You never add meta-commentary, "
            "titles-within-the-poem duplicated from the heading, author notes, or "
            "disclaimers. You output the poem text only, with its line breaks preserved."
        )
    if mode == "report":
        allowed_blocks = ", ".join(doc_profile.get("allowed_blocks", []))
        content_rules = "\n".join(f"- {rule}" for rule in doc_profile.get("content_rules", []))
        return (
            "You are a professional business-writing engine producing formal executive "
            f"{doc_profile.get('label', 'report').lower()} content. You write clear, precise, well-organized prose appropriate "
            "for senior stakeholders: short paragraphs, active voice, concrete claims, no "
            "fluff, no marketing language, no meta-commentary or disclaimers. Use structured JSON blocks for tables, charts, matrices, callouts, and other visual elements instead of raw Markdown when possible. "
            "When you include a data block, make it internally consistent and grounded in the brief and any knowledge base material provided; if real-world statistics are unavailable, label estimates clearly. Never invent unsupported facts. "
            f"Supported blocks: {allowed_blocks}.\n"
            f"Content rules:\n{content_rules}\n"
            "Use '## ' for a subheading only if the section genuinely needs internal subdivision. Never use tables or charts outside of the requested document structure."
        )
    # novel / short_story
    return (
        "You are a professional long-form fiction writing engine operating with full "
        "creative latitude, constrained only by the story context given to you. You write "
        "immersive, detailed, technically grounded prose — full scenes, dialogue, sensory "
        "detail, pacing, tension, and consequence. You never summarize when you should "
        "dramatize. You never break the fourth wall, add meta-commentary, headers, author "
        "notes, tables, charts, or disclaimers. You continue the manuscript seamlessly, in the tense and "
        "point of view already established, never repeating content already written."
    )


def writer_prompt(mode: str, proj: dict, outline_str: str, chapter: dict,
                   prior_context: str, unit_tail: str, words_so_far: int,
                   unit_target_words: int, first_pass: bool,
                   correction_notes: str = "", story_bible: str = "",
                   kb_context: str = "", doc_profile: Optional[dict] = None) -> str:
    m = MODES[mode]
    doc_profile = doc_profile or project_document_profile(proj)
    remaining = max(unit_target_words - words_so_far, 100)

    if mode == "poetry":
        closing_instruction = (
            "\nAim to complete the whole poem within this response. End your response "
            "with the exact token [UNIT_END] on its own final line once the poem is "
            "complete. If it genuinely needs one more short response to finish, omit "
            "the token."
        )
    elif remaining <= 350:
        closing_instruction = (
            "\nThis section/chapter is near its target length. Bring it to a natural, "
            "satisfying close within this response. End your response with the exact "
            "token [UNIT_END] on its own final line."
        )
    else:
        closing_instruction = (
            "\nDo NOT end yet — there is more content to cover here. If, and only if, "
            "it reaches a fully natural conclusion within this response anyway, end with "
            "[UNIT_END] on its own final line. Otherwise do not include that token."
        )

    if mode == "poetry":
        opening = "Write the poem now, in full." if first_pass else (
            "Continue directly from the lines above to finish the poem. Do not repeat "
            "lines already written, do not restate the title."
        )
    else:
        opening = "Begin now." if first_pass else (
            "Continue directly from the text above. Do not restate or repeat anything "
            "already written. Do not write a new title or heading."
        )

    correction_block = ""
    if correction_notes:
        correction_block = f"""
A reviewer flagged the following issues. Resolve them naturally as you continue — do
not call attention to the correction, just get it right from here on:
{correction_notes}
"""

    bible_block = f"\nSTORY BIBLE (established facts — stay consistent with these):\n{story_bible}\n" if story_bible else ""
    kb_block = f"\nKNOWLEDGE BASE — REFERENCE MATERIAL FOR THIS BOOK ONLY (ground relevant details in this where applicable):\n{kb_context}\n" if kb_context else ""
    style_block = style_notes_block(proj)
    chapter_meta = chapter.get("meta") or {}
    if isinstance(chapter_meta, str):
        try:
            chapter_meta = json.loads(chapter_meta)
        except Exception:
            chapter_meta = {}

    table_instruction = ""
    if mode == "report":
        if chapter_meta.get("include_table"):
            table_instruction = (f"\nThis section MUST include a structured JSON data block for a table. Table should contain: {chapter_meta.get('table_spec','')}\n")
        if chapter_meta.get("include_chart"):
            table_instruction += (f"\nThis section SHOULD include a structured JSON data block for a chart. Chart should contain: {chapter_meta.get('chart_spec','')}\n")
        table_instruction += "\nWhen data, comparisons, timelines, risk factors, statistics, or structured information appear anywhere, prefer structured JSON blocks over prose. Visual blocks are always preferred for quantitative or comparative content.\n"

    label = m["unit"]
    header = f'{"NOVEL" if mode=="novel" else "COLLECTION" if mode in ("short_story","poetry") else "REPORT"}: "{proj["title"]}"  |  {"GENRE" if mode!="report" else "DOMAIN"}: {proj["genre"]}'
    novel_meta = ""
    if mode in ("novel", "short_story") and chapter_meta:
        novel_meta = f"\nPOV CHARACTER: {chapter_meta.get('pov_character','')}\nPRIMARY SETTING: {chapter_meta.get('setting','')}\nTONE FOR THIS CHAPTER: {chapter_meta.get('tone','')}\n"
        if chapter_meta.get("epigraph"):
            novel_meta += "Do NOT restate the epigraph — it will be placed before your text by the typesetter.\n"

    return f"""{header}

PREMISE / BRIEF:
{proj['premise']}

FULL OUTLINE (for global consistency only — do not write other {m['unit_plural'].lower()}):
{outline_str}
{bible_block}{kb_block}{style_block}
RELEVANT CONTEXT RETRIEVED FROM EARLIER {m['unit_plural'].upper()}:
{prior_context if prior_context else f"(This is the opening {label.lower()} — no prior context.)"}

CURRENT {label.upper()} TO WRITE: "{chapter['title']}"
SYNOPSIS: {chapter['synopsis']}
{f"FORM: {chapter_meta.get('form','')}" if mode == "poetry" else f"KEY POINTS TO INCLUDE: {', '.join(chapter.get('key_events', []) or [])}"}
{novel_meta}
{f"DOCUMENT TYPE: {doc_profile.get('label', '')}" if mode == 'report' else ''}
{table_instruction}{correction_block}
TEXT WRITTEN SO FAR IN THIS {label.upper()} (tail shown for continuity):
{unit_tail if unit_tail else "(nothing written yet)"}

Target length: ~{unit_target_words} words. Written so far: ~{words_so_far} words.
{opening}{closing_instruction}
Write substantial, fully-developed content now.""".strip()


# ---- Continuity / Review -----------------------------------------------------

def continuity_sys_for_mode(mode: str) -> str:
    if mode == "report":
        return (
            "You are a report-review engine. You check a section draft against its brief "
            "and against the knowledge base / prior sections for factual inconsistency, "
            "unsupported specific claims, off-brand tone, or structural problems. You "
            "output a short, precise bullet list of concrete issues found, or the single "
            "word 'CLEAN' if there are none. You never rewrite prose yourself."
        )
    if mode == "novel":
        return (
            "CONTINUITY CHECKLIST — evaluate each item:\n"
            "1. Character names consistent with story bible?\n"
            "2. POV character correct for this chapter (matches outline)?\n"
            "3. Timeline coherent (no anachronisms vs. prior chapters)?\n"
            "4. No facts directly contradicted by story bible entities?\n"
            "5. Tone/register consistent with established narrative voice?\n"
            "6. No repeated plot points already resolved in prior chapters?\n"
            "7. Chapter ends at a dramatically appropriate beat (not mid-sentence or mid-scene unless cliffhanger)?\n\n"
            "If all 7 pass: respond with the single word PASS.\n"
            "If any fail: respond with 'REVISE:' followed by specific, actionable correction notes for the writer, referencing exact text if possible."
        )
    if mode == "poetry":
        return (
            "You are a poetry-review engine. You check a poem draft against its intended "
            "subject, form, and the collection's stated theme for genuine problems (broken "
            "requested form, drifted subject, inconsistent imagery/voice vs. the rest of "
            "the collection). You output a short, precise bullet list of concrete issues, "
            "or the single word 'CLEAN' if there are none. You never rewrite the poem "
            "yourself and never nitpick stylistic choices that are simply artistic."
        )
    return (
        "You are a continuity-and-consistency review engine. You check a draft against "
        "the outline and prior summaries for contradictions: character names, timelines, "
        "established facts, world rules, and tone. You output a short, precise bullet list "
        "of concrete issues found, or the single word 'CLEAN' if there are none. You never "
        "rewrite prose yourself and never invent issues that are not actually present."
    )


def continuity_prompt(mode: str, chapter: dict, chapter_text: str, prior_context: str,
                       kb_context: str = "") -> str:
    kb_block = f"\nKNOWLEDGE BASE MATERIAL TO CHECK AGAINST:\n{kb_context}\n" if kb_context else ""
    meta = chapter.get("meta") or {}
    if isinstance(meta, str):
        try:
            meta = json.loads(meta)
        except Exception:
            meta = {}
    extra = ""
    if mode == "novel":
        extra = f"\nPOV CHARACTER EXPECTED: {meta.get('pov_character','')}\nPRIMARY SETTING: {meta.get('setting','')}\n"
    return f"""OUTLINE FOR THIS ITEM: {chapter['title']} — {chapter['synopsis']}{extra}

RELEVANT PRIOR CONTEXT:
{prior_context if prior_context else "(none — opening item)"}
{kb_block}
DRAFT TO CHECK:
{chapter_text[:6500]}

List any genuine issues as short bullets. If there are no issues, respond with exactly:
CLEAN""".strip()


SUMMARY_SYS = (
    "You are a factual summarization engine for a writing system's long-term memory. You "
    "produce dense, precise summaries capturing only content-critical facts, states, and "
    "details needed for future items to remain consistent. No flourishes, no opinions, "
    "no filler."
)

def summary_prompt(chapter: dict, chapter_text: str) -> str:
    return f"""Summarize the following text in under 150 words. Capture: what happened or
was established, states/changes, and any facts future items must remain consistent with.

TITLE: {chapter['title']}
TEXT:
{chapter_text[:8000]}

Summary:""".strip()


# =============================================================================
# JSON PARSING HELPERS
# =============================================================================

def safe_json_array(text: str) -> list:
    cleaned = text.strip()
    cleaned = re.sub(r"^```(json)?", "", cleaned, flags=re.IGNORECASE).strip()
    cleaned = re.sub(r"```$", "", cleaned).strip()
    start = cleaned.find("[")
    end = cleaned.rfind("]")
    if start == -1 or end == -1 or end <= start:
        raise ValueError("No JSON array found in architect output.")
    snippet = cleaned[start:end + 1]
    return json.loads(snippet)


def architect_outline_with_retry(mode: str, architect: dict, title: str, genre: str,
                                  premise: str, num_units: int, words_per_unit: int,
                                  proj: dict, project_id: str, kb_context: str = "") -> list:
    """Ask the architect for an outline; if the JSON fails to parse, ask again
    (up to MAX_OUTLINE_PARSE_ATTEMPTS times) with an explicit correction notice
    instead of crashing the whole run over a formatting slip."""
    last_err = None
    sys_prompt = architect_sys_for_mode(mode, project_document_profile(proj))
    for attempt in range(1, MAX_OUTLINE_PARSE_ATTEMPTS + 1):
        prompt = architect_prompt(mode, title, genre, premise, num_units, words_per_unit,
                                   proj, kb_context)
        if attempt > 1:
            prompt += (f"\n\nIMPORTANT: Your previous response could not be parsed as JSON "
                       f"({last_err}). Output ONLY a valid JSON array this time — no prose, "
                       f"no markdown fences, no explanation.")
        raw = ollama_chat(architect["model"], sys_prompt, prompt,
                           temperature=0.85, num_predict=OUTLINE_TOKEN_TARGET)
        try:
            outline = safe_json_array(raw)
            if not isinstance(outline, list) or len(outline) == 0:
                raise ValueError("Parsed JSON was not a non-empty list.")
            cleaned = []
            for item in outline:
                if not isinstance(item, dict) or "title" not in item:
                    continue
                meta = {}
                if mode == "poetry":
                    meta["form"] = str(item.get("form", "")).strip()
                if mode == "report":
                    meta["include_table"] = bool(item.get("include_table", False))
                    meta["table_spec"] = str(item.get("table_spec", "")).strip()
                    meta["include_chart"] = bool(item.get("include_chart", False))
                    meta["chart_spec"] = str(item.get("chart_spec", "")).strip()
                cleaned.append({
                    "title": str(item.get("title", "Untitled")),
                    "synopsis": str(item.get("synopsis", "")),
                    "key_events": item.get("key_events", []) if isinstance(item.get("key_events"), list) else [],
                    "meta": meta,
                })
            if not cleaned:
                raise ValueError("No valid items found in outline.")
            return cleaned
        except Exception as e:
            last_err = str(e)
            log(project_id, architect["name"], "architect",
                f"Outline parse attempt {attempt} failed: {last_err}")
            if attempt < MAX_OUTLINE_PARSE_ATTEMPTS:
                time.sleep(2)
    raise OllamaError(f"Architect failed to produce a parseable outline after "
                       f"{MAX_OUTLINE_PARSE_ATTEMPTS} attempts: {last_err}")


def needs_rewrite(continuity_review: str) -> bool:
    if not continuity_review:
        return False
    text = continuity_review.strip().lower()
    if text == "clean" or text.startswith("clean"):
        return False
    return any(kw in text for kw in CONTINUITY_REWRITE_TRIGGER_WORDS) or len(text) > 20


# =============================================================================
# STRUCTURED CONTENT BLOCK PARSER
# Supports fenced JSON blocks for report-style documents, while preserving
# plain prose for literary modes. Report blocks can express tables, charts,
# matrices, callouts, and other visual assets as structured data rather than
# raw Markdown.
# =============================================================================

def _parse_pipe_table(raw: str) -> Optional[List[List[str]]]:
    lines = [l for l in raw.split("\n") if l.strip()]
    if len(lines) < 2:
        return None
    if not all("|" in l for l in lines[:2]) or not re.match(r"^\s*\|?\s*-+", lines[1]):
        return None
    rows = []
    for l in lines:
        if re.match(r"^\s*\|?\s*-+", l):
            continue
        cells = [c.strip() for c in l.strip().strip("|").split("|")]
        rows.append(cells)
    return rows if rows else None


def _parse_json_block(payload: Any) -> List[Tuple[str, Any]]:
    blocks: List[Tuple[str, Any]] = []
    if isinstance(payload, list):
        for item in payload:
            blocks.extend(_parse_json_block(item))
        return blocks
    if not isinstance(payload, dict):
        return [("para", str(payload))]

    block_type = str(payload.get("type") or payload.get("block_type") or payload.get("kind") or "").strip().lower()
    if not block_type:
        block_type = "para"

    if block_type == "table":
        headers = payload.get("headers") or payload.get("columns") or []
        rows = payload.get("rows") or payload.get("data") or []
        if headers and rows:
            return [("table", [list(headers)] + [list(r) if isinstance(r, (list, tuple)) else [str(r)] for r in rows])]
        if rows and all(isinstance(r, (list, tuple)) for r in rows):
            return [("table", [list(r) for r in rows])]

    if block_type in ("chart", "bar_chart", "line_chart", "pie_chart", "donut_chart", "scatter_chart", "timeline_chart", "gantt_chart", "heatmap", "risk_matrix", "radar_chart", "funnel_chart", "treemap", "timeline", "kpi_cards"):
        chart_spec = dict(payload)
        chart_spec["chart_type"] = payload.get("chart_type") or block_type
        return [("chart", chart_spec)]

    if block_type in ("callout_box", "note_box", "warning_box", "tip_box", "deliverable_cards", "persona_cards", "severity_badge", "analyst_notes", "executive_summary_block", "ioc_table", "cve_table", "mitre_table", "sigma_rule_block", "yara_rule_block", "network_diagram", "stix_object_table", "action_items_table", "remediation_table", "lessons_learned_block"):
        return [(block_type, payload)]

    if block_type in ("appendix", "glossary", "bibliography", "figure", "image_caption", "equation", "footnote", "section_divider", "page_break", "heading", "quote", "bullet_list", "numbered_list", "paragraph", "code_block", "command_block", "timeline", "risk_matrix", "swot_matrix", "comparison_matrix", "parameter_table", "api_reference", "troubleshooting", "faq", "changelog", "signature_page", "source_notes"):
        return [(block_type, payload)]

    return [("para", payload.get("text") or payload.get("content") or json.dumps(payload, ensure_ascii=False))]


def parse_content_blocks(text: str) -> List[Tuple[str, Any]]:
    blocks: List[Tuple[str, Any]] = []
    pattern = re.compile(r"```([a-zA-Z0-9_-]+)?\s*\n(.*?)\n```", re.DOTALL)
    pos = 0
    for match in pattern.finditer(text or ""):
        prefix = text[pos:match.start()]
        blocks.extend(_parse_plain_blocks(prefix))
        lang = (match.group(1) or "").strip().lower()
        code = match.group(2).strip()
        if lang == "json" or code.startswith("{") or code.startswith("["):
            try:
                blocks.extend(_parse_json_block(json.loads(code)))
            except Exception:
                blocks.append(("code_block", {"text": code, "language": lang or "text"}))
        else:
            blocks.append(("code_block", {"text": code, "language": lang or "text"}))
        pos = match.end()
    blocks.extend(_parse_plain_blocks(text[pos:]))
    return [b for b in blocks if b and str(b[1]).strip()]


def _parse_plain_blocks(text: str) -> List[Tuple[str, Any]]:
    blocks: List[Tuple[str, Any]] = []
    raw_blocks = [b for b in re.split(r"\n\s*\n", text or "") if b.strip()]
    for raw in raw_blocks:
        stripped = raw.strip()
        lines = [l for l in raw.split("\n") if l.strip()]
        table_rows = _parse_pipe_table(raw)
        if table_rows:
            blocks.append(("table", table_rows))
            continue
        if all(re.match(r"^\s*([-*]|\d+[.)])\s+", l) for l in lines):
            items = [re.sub(r"^\s*([-*]|\d+[.)])\s+", "", l).strip() for l in lines]
            blocks.append(("bullet_list", items))
            continue
        if stripped.startswith("## "):
            blocks.append(("heading", stripped[3:].strip()))
            continue
        if stripped.startswith(">"):
            blocks.append(("quote", stripped.lstrip("> ").strip()))
            continue
        blocks.append(("paragraph", stripped))
    return blocks


def _clean_text(value: Any) -> str:
    if isinstance(value, str):
        return value.strip()
    return str(value or "").strip()


def _normalise_list(value: Any) -> List[str]:
    if isinstance(value, list):
        return [_clean_text(v) for v in value if _clean_text(v)]
    if value is None:
        return []
    return [_clean_text(value)]


def render_callout_block(payload: dict, styles: Dict[str, ParagraphStyle], theme: dict, kind: str = "note") -> Table:
    title = _clean_text(payload.get("title") or kind.title())
    body = _clean_text(payload.get("text") or payload.get("content") or payload.get("body"))
    accent_map = {
        "note": rl_colors.HexColor("#1e6fa5"),
        "warning": rl_colors.HexColor("#d97706"),
        "tip": rl_colors.HexColor("#16a34a"),
        "danger": rl_colors.HexColor("#dc2626"),
        "finding": theme["pdf_accent"],
        "analyst_note": rl_colors.HexColor("#7c3aed"),
    }
    left_color = accent_map.get(kind, theme["pdf_accent"])
    title_style = ParagraphStyle(f"{kind}_title", fontName=styles["subheading"].fontName, fontSize=10, leading=12,
                                 textColor=left_color, spaceAfter=3)
    body_style = ParagraphStyle(f"{kind}_body", parent=styles["body"], italic=(kind in ("analyst_note", "finding")))
    box = Table([[Paragraph(f"<b>{escape_html(title)}</b>", title_style), Paragraph(escape_html(body), body_style)]], colWidths=[1.3*inch, 5.0*inch])
    box.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, -1), rl_colors.HexColor("#fffdf6")),
        ("LINEBEFORE", (0, 0), (0, -1), 4, left_color),
        ("BOX", (0, 0), (-1, -1), 0.6, theme["pdf_muted"]),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("LEFTPADDING", (0, 0), (-1, -1), 6),
        ("RIGHTPADDING", (0, 0), (-1, -1), 6),
        ("TOPPADDING", (0, 0), (-1, -1), 5),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
    ]))
    return box


def render_severity_badge(level: str, styles: Dict[str, ParagraphStyle]) -> Table:
    palette = {
        "critical": rl_colors.HexColor("#d32f2f"),
        "high": rl_colors.HexColor("#e65100"),
        "medium": rl_colors.HexColor("#f9a825"),
        "low": rl_colors.HexColor("#0277bd"),
        "info": rl_colors.HexColor("#546e7a"),
    }
    color = palette.get(level.lower(), rl_colors.HexColor("#546e7a"))
    badge = Table([[Paragraph(f"<b>{escape_html(level.upper())}</b>", styles["bullet"]) ]], colWidths=[1.1 * inch])
    badge.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, -1), color),
        ("TEXTCOLOR", (0, 0), (-1, -1), rl_colors.white),
        ("ALIGN", (0, 0), (-1, -1), "CENTER"),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("LEFTPADDING", (0, 0), (-1, -1), 4),
        ("RIGHTPADDING", (0, 0), (-1, -1), 4),
        ("TOPPADDING", (0, 0), (-1, -1), 2),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 2),
    ]))
    return badge


def render_kpi_cards(payload: dict, styles: Dict[str, ParagraphStyle], theme: dict) -> Table:
    cards = payload.get("cards") or []
    cells = []
    for card in cards[:4]:
        label = _clean_text(card.get("label"))
        value = _clean_text(card.get("value"))
        delta = _clean_text(card.get("delta"))
        delta_positive = bool(card.get("delta_positive", True))
        delta_color = rl_colors.HexColor("#16a34a" if delta_positive else "#dc2626")
        html = f"<font size='8'><b>{escape_html(label.upper())}</b></font><br/><font size='16'><b>{escape_html(value)}</b></font><br/><font size='8' color='{delta_color}'>{escape_html(delta)}</font>"
        cells.append(Paragraph(html, styles["body"]))
    if not cells:
        cells = [Paragraph("No KPI data", styles["body"])]
    table = Table([cells], colWidths=[6.2 * inch / max(len(cells), 1)] * len(cells))
    table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, -1), rl_colors.white),
        ("BOX", (0, 0), (-1, -1), 0.6, theme["pdf_muted"]),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("LEFTPADDING", (0, 0), (-1, -1), 6),
        ("RIGHTPADDING", (0, 0), (-1, -1), 6),
        ("TOPPADDING", (0, 0), (-1, -1), 6),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
    ]))
    return table


def render_timeline_block(payload: dict, styles: Dict[str, ParagraphStyle], theme: dict) -> Table:
    events = payload.get("events") or []
    rows = []
    for event in events:
        date = Paragraph(f"<font face='{styles['mono'].fontName}' size='8'><b>{escape_html(_clean_text(event.get('date')))}</b></font>", styles["body"])
        title = f"<b>{escape_html(_clean_text(event.get('title')))}</b><br/><i>{escape_html(_clean_text(event.get('detail') or event.get('description') or ''))}</i>"
        rows.append([date, Paragraph(title, styles["body"])])
    if not rows:
        rows = [[Paragraph("", styles["body"]), Paragraph("No timeline data", styles["body"])]]
    table = Table(rows, colWidths=[1.1 * inch, 5.1 * inch])
    table.setStyle(TableStyle([
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("BOX", (0, 0), (-1, -1), 0.6, theme["pdf_muted"]),
        ("LINEBEFORE", (0, 0), (0, -1), 2, theme["pdf_accent"]),
        ("LEFTPADDING", (0, 0), (-1, -1), 6),
        ("RIGHTPADDING", (0, 0), (-1, -1), 6),
        ("TOPPADDING", (0, 0), (-1, -1), 5),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
    ]))
    return table


def render_ioc_table(payload: dict, styles: Dict[str, ParagraphStyle], theme: dict) -> Table:
    rows = payload.get("rows") or payload.get("data") or []
    table_rows = [["Indicator", "Type", "Confidence", "First Seen", "Context"]]
    for row in rows:
        table_rows.append([
            _clean_text(row.get("indicator") or row.get("ioc") or row.get("value")),
            _clean_text(row.get("type") or row.get("kind")),
            _clean_text(row.get("confidence") or row.get("score")),
            _clean_text(row.get("first_seen") or row.get("date")),
            _clean_text(row.get("context") or row.get("notes")),
        ])
    return render_table_block(table_rows, theme, styles)


def render_code_block(payload: dict, styles: Dict[str, ParagraphStyle], theme: dict, header: str = "") -> Table:
    text = _clean_text(payload.get("text") or payload.get("content") or payload.get("code") or payload.get("body"))
    body = Paragraph(f"<font face='{styles['mono'].fontName}' size='8'>{escape_html(text).replace(chr(10), '<br/>')}</font>", styles["body"])
    title = Paragraph(f"<b>{escape_html(header)}</b>", styles["subheading"]) if header else None
    rows = [[title, body]] if title else [[body]]
    table = Table(rows, colWidths=[6.2 * inch])
    table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, -1), rl_colors.HexColor("#f5f5f5")),
        ("BOX", (0, 0), (-1, -1), 0.6, theme["pdf_muted"]),
        ("LEFTPADDING", (0, 0), (-1, -1), 6),
        ("RIGHTPADDING", (0, 0), (-1, -1), 6),
        ("TOPPADDING", (0, 0), (-1, -1), 6),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
    ]))
    return table


def render_fallback_block(block_type: str, payload: Any, styles: Dict[str, ParagraphStyle]) -> Paragraph:
    if isinstance(payload, dict):
        content = payload.get("content") or payload.get("text") or payload.get("body") or json.dumps(payload, ensure_ascii=False)
    else:
        content = str(payload)
    return Paragraph(escape_html(f"[BLOCK: {block_type}]\n{content}"), styles["body"])


def render_block_for_pdf(block_type: str, payload: Any, story: List[Any], styles: Dict[str, ParagraphStyle], theme: dict, profile: Dict[str, Any], counters: Dict[str, Any]) -> None:
    allowed = set(profile.get("allowed_blocks", []))
    disallowed = set(profile.get("disallowed_blocks", []))
    if block_type in disallowed:
        logger.warning("Skipping disallowed block '%s' for %s", block_type, profile.get("document_type"))
        return
    if allowed and block_type not in allowed and block_type not in {"paragraph", "heading", "subheading", "quote", "bullet_list", "numbered_list", "page_break", "section_divider"}:
        story.append(render_fallback_block(block_type, payload, styles))
        return

    if block_type == "table":
        story.append(render_table_block(payload if isinstance(payload, list) else [[str(payload)]], theme, styles))
    elif block_type in ("chart", "bar_chart", "line_chart", "pie_chart", "donut_chart", "scatter_chart", "timeline_chart", "gantt_chart", "heatmap", "risk_matrix", "swot_matrix", "comparison_matrix", "timeline"):
        story.append(render_chart_block(payload if isinstance(payload, dict) else {"type": block_type, "data": payload}, theme))
    elif block_type == "kpi_cards":
        story.append(render_kpi_cards(payload if isinstance(payload, dict) else {}, styles, theme))
    elif block_type in ("callout_box", "note_box", "warning_box", "tip_box", "finding", "analyst_note", "analyst_notes", "executive_summary_block"):
        kind = block_type.replace("_box", "")
        if kind == "analyst_notes":
            kind = "analyst_note"
        story.append(render_callout_block(payload if isinstance(payload, dict) else {"text": str(payload)}, styles, theme, kind=kind))
    elif block_type in ("sigma_rule_block", "yara_rule_block"):
        header = "SIGMA RULE" if block_type == "sigma_rule_block" else "YARA RULE"
        story.append(render_code_block(payload if isinstance(payload, dict) else {"text": str(payload)}, styles, theme, header=header))
    elif block_type == "code_block":
        story.append(render_code_block(payload if isinstance(payload, dict) else {"text": str(payload)}, styles, theme))
    elif block_type == "severity_badge":
        level = _clean_text(payload.get("level") if isinstance(payload, dict) else payload)
        story.append(render_severity_badge(level, styles))
    elif block_type == "ioc_table":
        story.append(render_ioc_table(payload if isinstance(payload, dict) else {}, styles, theme))
    elif block_type == "timeline":
        story.append(render_timeline_block(payload if isinstance(payload, dict) else {}, styles, theme))
    elif block_type == "image_caption":
        counters["figure"] = counters.get("figure", 0) + 1
        cap = _clean_text(payload.get("text") if isinstance(payload, dict) else payload)
        story.append(Paragraph(f"<i>Figure {counters['figure']}. {escape_html(cap)}</i>", ParagraphStyle("caption", parent=styles["body"], fontSize=8, leading=10, alignment=TA_CENTER, spaceBefore=2, spaceAfter=8)))
    elif block_type == "figure":
        story.append(render_callout_block(payload if isinstance(payload, dict) else {"text": str(payload)}, styles, theme, kind="note"))
    elif block_type == "page_break":
        story.append(PageBreak())
    elif block_type == "section_divider":
        story.append(Spacer(1, 12)); story.append(Paragraph("* * *", styles["chapter_num"])); story.append(Spacer(1, 12))
    elif block_type == "appendix":
        story.append(Paragraph("Appendix", styles["subheading"]))
        text_value = _clean_text(payload.get("text") if isinstance(payload, dict) else payload)
        if text_value:
            story.append(Paragraph(escape_html(text_value), styles["body"]))
    elif block_type in ("paragraph", "heading", "subheading", "quote"):
        text_value = _clean_text(payload)
        if block_type == "heading":
            counters["section"] = counters.get("section", 0) + 1
            if profile.get("section_numbering"):
                text_value = f"{counters['section']}. {text_value}"
            story.append(Paragraph(escape_html(text_value), styles["subheading"]))
        elif block_type == "quote":
            story.append(Paragraph(f"<i>{escape_html(text_value)}</i>", styles["body"]))
        else:
            story.append(Paragraph(escape_html(text_value), styles["body"]))
    elif block_type in ("bullet_list", "numbered_list"):
        items = payload if isinstance(payload, list) else _normalise_list(payload)
        story.append(ListFlowable([ListItem(Paragraph(escape_html(it), styles["bullet"]), leftIndent=12) for it in items], bulletType="bullet", start="•"))
    else:
        story.append(render_fallback_block(block_type, payload, styles))


def render_table_block(rows: List[List[str]], theme: dict, styles: Dict[str, ParagraphStyle]) -> Table:
    if not rows:
        rows = [["(empty)"]]
    header = rows[0]
    body_rows = rows[1:] if len(rows) > 1 else []
    header_cells = [Paragraph(escape_html(h), styles["table_header"]) for h in header]
    data = [header_cells]
    for r in body_rows:
        r = r + [""] * (len(header) - len(r))
        data.append([Paragraph(escape_html(c), styles["table_cell"]) for c in r[:len(header)]])
    n_cols = max(len(header), 1)
    col_width = 6.3 * inch / n_cols
    t = Table(data, colWidths=[col_width] * n_cols, repeatRows=1)
    header_bg = theme.get("pdf_table_header_bg", theme["pdf_accent"])
    row_alt = theme.get("pdf_table_row_alt", rl_colors.HexColor("#f2f2f2"))
    style_cmds = [
        ("BACKGROUND", (0, 0), (-1, 0), header_bg),
        ("GRID", (0, 0), (-1, -1), 0.5, theme["pdf_muted"]),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("TOPPADDING", (0, 0), (-1, -1), 5),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
    ]
    for i in range(1, len(data)):
        if i % 2 == 0:
            style_cmds.append(("BACKGROUND", (0, i), (-1, i), row_alt))
    t.setStyle(TableStyle(style_cmds))
    return t


def render_chart_block(spec: dict, theme: dict) -> Drawing:
    chart_type = str(spec.get("chart_type") or spec.get("type") or "chart").lower()
    title = _clean_text(spec.get("title") or spec.get("label") or spec.get("name") or chart_type.replace("_", " ").title())
    width = 6.3 * inch
    height = 3.2 * inch
    drawing = Drawing(width, height)
    drawing.add(Rect(0, 0, width, height, fillColor=rl_colors.white, strokeColor=theme["pdf_muted"], strokeWidth=0.6))
    drawing.add(String(12, height - 16, title, fontName=theme["pdf_head_font"], fontSize=10.5, fillColor=theme["pdf_ink"]))

    if chart_type in ("bar", "bar_chart", "column", "column_chart"):
        values = spec.get("values") or spec.get("data") or []
        if values and isinstance(values[0], dict):
            labels = [_clean_text(v.get("label") or v.get("name")) for v in values]
            series = [float(v.get("value", 0) or 0) for v in values]
        else:
            labels = _normalise_list(spec.get("labels"))
            series = [float(v) for v in values] if isinstance(values, list) else []
        chart = VerticalBarChart()
        chart.x = 40
        chart.y = 40
        chart.height = 200
        chart.width = width - 90
        chart.data = [series or [0]]
        chart.categoryAxis.categoryNames = labels or [str(i + 1) for i in range(len(series or [0]))]
        chart.bars[0].fillColor = theme["pdf_accent"]
        chart.valueAxis.valueMin = 0
        drawing.add(chart)
        return drawing

    if chart_type in ("line", "line_chart", "timeline_chart", "gantt_chart"):
        series = spec.get("series") or spec.get("data") or []
        plot = LinePlot()
        plot.x = 40
        plot.y = 45
        plot.height = 190
        plot.width = width - 90
        plot.data = []
        if series and isinstance(series, list) and isinstance(series[0], dict):
            for serie in series:
                points = []
                for point in serie.get("points", []):
                    if isinstance(point, (list, tuple)) and len(point) >= 2:
                        points.append((float(point[0]), float(point[1])))
                if points:
                    plot.data.append(points)
        elif series and isinstance(series, list):
            points = []
            for idx, value in enumerate(series):
                points.append((idx + 1, float(value)))
            plot.data.append(points)
        if not plot.data:
            plot.data = [[(1, 0), (2, 1)]]
        plot.lines[0].strokeColor = theme["pdf_accent"]
        plot.lines[0].strokeWidth = 2
        drawing.add(plot)
        return drawing

    if chart_type in ("pie", "pie_chart", "donut", "donut_chart"):
        values = spec.get("values") or spec.get("data") or []
        labels = _normalise_list(spec.get("labels"))
        if not values and isinstance(spec.get("slices"), list):
            values = [float(v.get("value", 0) or 0) for v in spec.get("slices", [])]
            labels = [_clean_text(v.get("label") or v.get("name")) for v in spec.get("slices", [])]
        pie = Pie()
        pie.x = 120
        pie.y = 28
        pie.width = 160
        pie.height = 160
        pie.data = [float(v) for v in values] if values else [1, 1]
        pie.labels = labels or [str(i + 1) for i in range(len(pie.data))]
        pie.slices.strokeWidth = 0.5
        if chart_type in ("donut", "donut_chart"):
            pie.innerRadiusFraction = 0.55
        drawing.add(pie)
        return drawing

    if chart_type in ("scatter", "scatter_chart"):
        points = spec.get("points") or spec.get("data") or []
        plot = ScatterPlot()
        plot.x = 40
        plot.y = 45
        plot.height = 190
        plot.width = width - 90
        series_points = []
        for point in points:
            if isinstance(point, (list, tuple)) and len(point) >= 2:
                series_points.append((float(point[0]), float(point[1])))
        plot.data = [series_points or [(1, 1), (2, 2)]]
        plot.symbol = Circle(0, 0, 3)
        plot.symbol.fillColor = theme["pdf_accent"]
        plot.symbol.strokeColor = theme["pdf_accent"]
        drawing.add(plot)
        return drawing

    if chart_type in ("risk_matrix", "heatmap", "swot_matrix", "comparison_matrix"):
        matrix = spec.get("matrix") or spec.get("cells") or []
        x_labels = _normalise_list(spec.get("x_labels") or spec.get("columns"))
        y_labels = _normalise_list(spec.get("y_labels") or spec.get("rows"))
        cols = max(len(x_labels), len(matrix[0]) if matrix and isinstance(matrix[0], list) else 5)
        rows = max(len(y_labels), len(matrix) if matrix else 5)
        cell_w = 380 / max(cols, 1)
        cell_h = 170 / max(rows, 1)
        origin_x = 85
        origin_y = 35
        palette = [rl_colors.HexColor("#f4e3e3"), rl_colors.HexColor("#f7d7b5"), rl_colors.HexColor("#f6e2c4"), rl_colors.HexColor("#dbeacf"), rl_colors.HexColor("#bfd8bf")]
        for r in range(rows):
            for c in range(cols):
                value = 2
                if matrix and r < len(matrix) and isinstance(matrix[r], list) and c < len(matrix[r]):
                    try:
                        value = int(float(matrix[r][c]))
                    except Exception:
                        value = 2
                fill = palette[min(max(value, 0), len(palette) - 1)]
                x = origin_x + c * cell_w
                y = origin_y + (rows - 1 - r) * cell_h
                drawing.add(Rect(x, y, cell_w, cell_h, fillColor=fill, strokeColor=theme["pdf_muted"], strokeWidth=0.4))
        for idx, label in enumerate(x_labels[:cols]):
            drawing.add(String(origin_x + idx * cell_w + cell_w / 2, origin_y + rows * cell_h + 8, label, fontSize=7, textAnchor="middle", fillColor=theme["pdf_ink"]))
        for idx, label in enumerate(y_labels[:rows]):
            drawing.add(String(8, origin_y + (rows - 1 - idx) * cell_h + cell_h / 2, label, fontSize=7, fillColor=theme["pdf_ink"]))
        return drawing

    note = _clean_text(spec.get("text") or spec.get("content") or spec.get("description") or json.dumps(spec, ensure_ascii=False))
    drawing.add(String(16, 120, note[:120], fontName=theme["pdf_body_font"], fontSize=9, fillColor=theme["pdf_ink"]))
    return drawing


def render_text_block(text: str, styles: Dict[str, ParagraphStyle]) -> Paragraph:
    return Paragraph(escape_html(text), styles["body"])


# =============================================================================
# ROUND ROBIN AGENT SELECTION
# =============================================================================

class RoundRobin:
    def __init__(self, items: List[dict]):
        self.items = items
        self.i = 0

    def next(self) -> Optional[dict]:
        if not self.items:
            return None
        item = self.items[self.i % len(self.items)]
        self.i += 1
        return item


# =============================================================================
# GENERATION PIPELINE
# =============================================================================

def write_unit_body(project_id: str, proj: dict, outline_str: str, chapter: dict,
                     writer: dict, prior_context: str, unit_target_words: int,
                     existing_content: str = "", kb_context: str = "") -> str:
    """Runs the bounded multi-pass writing loop for a single unit (chapter /
    story / poem / report section). Resumable: pass existing_content to
    continue rather than restart that unit from scratch."""
    mode = proj["mode"]
    doc_profile = project_document_profile(proj)
    sys_prompt = writer_sys_for_mode(mode, doc_profile)
    unit_text = existing_content or ""
    passes = 0
    max_passes = 4 if mode == "poetry" else MAX_PASSES_PER_UNIT
    while word_count(unit_text) < unit_target_words and passes < max_passes:
        if STOP_FLAGS.get(project_id):
            break
        tail = " ".join(unit_text.split()[-350:]) if unit_text else ""
        prompt = writer_prompt(
            mode, proj, outline_str, chapter, prior_context, tail,
            word_count(unit_text), unit_target_words,
            first_pass=(passes == 0 and not existing_content),
            kb_context=kb_context,
            doc_profile=doc_profile,
        )
        temperature = 0.85 if mode == "report" else 0.93
        try:
            resp = ollama_chat(writer["model"], sys_prompt, prompt, temperature=temperature)
        except OllamaError as e:
            log(project_id, writer["name"], "error", f"Writer call failed: {e}")
            raise
        ended_naturally = "[UNIT_END]" in resp
        resp = resp.replace("[UNIT_END]", "").strip()
        unit_text = (unit_text + "\n\n" + resp).strip() if unit_text else resp
        passes += 1

        save_chapter_content(project_id, chapter["idx"], unit_text)
        done_words = total_words_written(project_id)
        elapsed = time.time() - PROGRESS.get(project_id, {}).get("_t_start", time.time())
        rate = done_words / elapsed if elapsed > 0 else 0
        remaining_words = max(proj["target_words"] - done_words, 0)
        eta = int(remaining_words / rate) if rate > 0 else 0
        update_progress(project_id, words_written=done_words, eta_seconds=eta)

        if ended_naturally:
            break
    return unit_text


def run_generation(project_id: str) -> None:
    lock = project_lock(project_id)
    if not lock.acquire(blocking=False):
        logger.info("Generation already running for %s — ignoring duplicate start.", project_id)
        return
    try:
        proj = get_project(project_id)
        if not proj:
            logger.error("run_generation called for unknown project %s", project_id)
            return
        mode = proj["mode"]
        m = MODES.get(mode, MODES["novel"])

        agents = get_agents(project_id)
        architects = [a for a in agents if a["role"] == "architect"]
        writers = [a for a in agents if a["role"] == "writer"]
        continuity_agents = [a for a in agents if a["role"] == "continuity"]

        if not architects or not writers:
            update_progress(project_id, status="error",
                             error="Need at least one architect and one writer agent.")
            return

        ok, msg = ollama_health()
        if not ok:
            update_progress(project_id, status="error", error=msg)
            log(project_id, "system", "error", msg)
            return

        available = set(ollama_list_models())
        if available:
            missing = [a["model"] for a in agents if a["model"] not in available]
            if missing:
                err = f"Models not found in Ollama: {', '.join(sorted(set(missing)))}. Pull them with 'ollama pull <model>'."
                update_progress(project_id, status="error", error=err)
                log(project_id, "system", "error", err)
                return

        architect = architects[0]
        writer_rr = RoundRobin(writers)
        continuity_rr = RoundRobin(continuity_agents)
        embed_model = proj["embed_model"]

        with PROGRESS_LOCK:
            PROGRESS.setdefault(project_id, {})["_t_start"] = time.time()

        existing_chapters = get_chapters(project_id)

        if existing_chapters:
            # ---------- RESUME PATH ----------
            outline = []
            for c in existing_chapters:
                try:
                    meta = json.loads(c["meta"] or "{}")
                except Exception:
                    meta = {}
                outline.append({"idx": c["idx"], "title": c["title"], "synopsis": c["synopsis"],
                                 "key_events": json.loads(c["key_events"] or "[]"), "meta": meta})
            num_units = len(outline)
            update_progress(project_id, status="running", phase="writing",
                             total_units=num_units, error=None)
            log(project_id, "system", "system",
                f"Resuming project: {sum(1 for c in existing_chapters if c['status']=='done')}/"
                f"{num_units} {m['unit_plural'].lower()} already complete.")
        else:
            # ---------- FRESH PLAN PATH ----------
            update_progress(project_id, status="running", phase="planning", error=None)
            num_units = proj["num_units"]
            words_per_unit = proj["words_per_unit"]
            log(project_id, architect["name"], "architect",
                f"Planning {num_units} {m['unit_plural'].lower()} (~{words_per_unit} words each).")

            kb_context = knowledge_context(project_id, proj["premise"], embed_model, top_k=8)

            raw_outline = architect_outline_with_retry(
                mode, architect, proj["title"], proj["genre"], proj["premise"],
                num_units, words_per_unit, proj, project_id, kb_context
            )
            num_units = len(raw_outline)

            with db() as conn:
                for idx, ch in enumerate(raw_outline):
                    conn.execute(
                        "INSERT INTO chapters (id, project_id, idx, title, synopsis, key_events, meta, status) "
                        "VALUES (?,?,?,?,?,?,?,?)",
                        (str(uuid.uuid4()), project_id, idx, ch["title"], ch["synopsis"],
                         json.dumps(ch["key_events"]), json.dumps(ch["meta"]), "pending")
                    )

            outline = [{**ch, "idx": i} for i, ch in enumerate(raw_outline)]
            update_progress(project_id, total_units=num_units, phase="writing")
            log(project_id, architect["name"], "architect",
                f"Outline complete: {num_units} {m['unit_plural'].lower()}.")

        outline_str = "\n".join(f"{c['idx']+1}. {c['title']} — {c['synopsis']}" for c in outline)
        unit_target_words = proj["words_per_unit"]

        for chapter in outline:
            idx = chapter["idx"]
            if STOP_FLAGS.get(project_id):
                update_progress(project_id, status="stopped", phase="stopped")
                log(project_id, "system", "system", "Generation stopped by user.")
                return

            row = get_chapter(project_id, idx)
            if row and row["status"] == "done":
                continue

            writer = writer_rr.next()
            cont_agent = continuity_rr.next()

            with db() as conn:
                conn.execute(
                    "UPDATE chapters SET status='writing', agent_name=? WHERE project_id=? AND idx=?",
                    (writer["name"], project_id, idx)
                )
            log(project_id, architect["name"], "architect",
                f"Brief for {m['unit']} {idx+1} ('{chapter['title']}') assigned to {writer['name']}.")

            query = f"{chapter.get('title','')} {chapter.get('synopsis','')}"
            retrieved = vector_retrieve(project_id, query, embed_model, top_k=6,
                                         kind="chapter_summary") if idx > 0 else []
            prior_context = "\n".join(f"- ({r[2]} / unit {r[1]+1}) {r[3]}" for r in retrieved)
            kb_context = knowledge_context(project_id, query, embed_model, top_k=5)

            story_bible = story_bible_context(project_id) if m["uses_story_bible"] else ""

            existing_content = row["content"] if row and row["status"] == "writing" else ""
            if existing_content:
                log(project_id, writer["name"], "writer",
                    f"Resuming {m['unit']} {idx+1} from {word_count(existing_content)} existing words.")

            try:
                unit_text = write_unit_body(
                    project_id, proj, outline_str, chapter, writer, prior_context,
                    unit_target_words, existing_content, kb_context
                )
            except OllamaError as e:
                update_progress(project_id, status="error", error=str(e))
                with db() as conn:
                    conn.execute("UPDATE chapters SET status='error' WHERE project_id=? AND idx=?",
                                 (project_id, idx))
                return

            if STOP_FLAGS.get(project_id):
                update_progress(project_id, status="stopped", phase="stopped")
                log(project_id, "system", "system", "Generation stopped by user mid-unit (progress saved).")
                return

            log(project_id, writer["name"], "writer",
                f"{m['unit']} {idx+1} draft complete ({word_count(unit_text)} words).")

            # ---- Review with one bounded corrective pass ----
            continuity_notes = ""
            if cont_agent:
                try:
                    review = ollama_chat(
                        cont_agent["model"], continuity_sys_for_mode(mode),
                        continuity_prompt(mode, chapter, unit_text, prior_context, kb_context),
                        temperature=0.3, num_predict=400
                    )
                except OllamaError as e:
                    review = "CLEAN"
                    log(project_id, cont_agent["name"], "error",
                        f"Review failed, skipping: {e}")
                log(project_id, cont_agent["name"], "continuity", review[:1500])
                continuity_notes = review

                if needs_rewrite(review):
                    log(project_id, writer["name"], "writer",
                        f"Applying one corrective pass for {m['unit']} {idx+1} based on review notes.")
                    try:
                        fix_prompt = writer_prompt(
                            mode, proj, outline_str, chapter, prior_context,
                            " ".join(unit_text.split()[-350:]),
                            word_count(unit_text), word_count(unit_text) + 250,
                            first_pass=False, correction_notes=review,
                            story_bible=story_bible, kb_context=kb_context,
                        )
                        fix_resp = ollama_chat(writer["model"], writer_sys_for_mode(mode),
                                                fix_prompt, temperature=0.7)
                        fix_resp = fix_resp.replace("[UNIT_END]", "").strip()
                        unit_text = (unit_text + "\n\n" + fix_resp).strip()
                        save_chapter_content(project_id, idx, unit_text)
                        with db() as conn:
                            conn.execute(
                                "UPDATE chapters SET revision_count = revision_count + 1 WHERE project_id=? AND idx=?",
                                (project_id, idx)
                            )
                    except OllamaError as e:
                        log(project_id, writer["name"], "error", f"Corrective pass failed, keeping draft as-is: {e}")

            # ---- Story bible update (novel / short_story only) ----
            if m["uses_story_bible"]:
                try:
                    extract_and_merge_entities(project_id, chapter, unit_text, architect["model"])
                except Exception as e:
                    log(project_id, "system", "error", f"Story bible update failed: {e}")

            # ---- Summarize for long-term memory ----
            try:
                summary = ollama_chat(
                    architect["model"], SUMMARY_SYS, summary_prompt(chapter, unit_text),
                    temperature=0.4, num_predict=300
                )
            except OllamaError as e:
                summary = unit_text[:600]
                log(project_id, architect["name"], "error",
                    f"Summary generation failed, falling back to raw excerpt: {e}")

            with db() as conn:
                conn.execute(
                    "UPDATE chapters SET status='done', summary=?, continuity_notes=? "
                    "WHERE project_id=? AND idx=?",
                    (summary, continuity_notes, project_id, idx)
                )

            vector_add(project_id, idx, "chapter_summary",
                       f"{m['unit']} {idx+1} ({chapter['title']}): {summary}", embed_model)

            units_done = sum(1 for c in get_chapters(project_id) if c["status"] == "done")
            update_progress(project_id, units_done=units_done)
            log(project_id, "system", "system", f"{m['unit']} {idx+1}/{num_units} complete.")

        update_progress(project_id, phase="compiling_pdf")
        log(project_id, "system", "system", "All units complete. Compiling PDF.")
        pdf_path = build_pdf(project_id)
        update_progress(project_id, status="completed", phase="completed")
        with db() as conn:
            conn.execute("UPDATE projects SET completed_at=? WHERE id=?",
                         (datetime.utcnow().isoformat(), project_id))
        log(project_id, "system", "system", f"Book complete. PDF: {pdf_path}")

    except Exception as e:
        tb = traceback.format_exc()
        update_progress(project_id, status="error", error=str(e))
        log(project_id, "system", "error", f"{e}\n{tb[-1500:]}")
        logger.exception("Generation crashed for project %s", project_id)
    finally:
        lock.release()


# =============================================================================
# PDF GENERATION — mode + theme aware. Tables/report banners only ever
# render in 'report' mode; novel/short_story/poetry stay pure literary
# typesetting regardless of which theme is selected.
# =============================================================================

def build_pdf_styles(theme: dict, mode: str, doc_profile: Optional[dict] = None) -> Dict[str, ParagraphStyle]:
    literary = MODES[mode]["literary"]
    profile = doc_profile or get_document_profile(mode)
    body_font = resolve_pdf_font(profile["fonts"]["body"], "body")
    heading_font = resolve_pdf_font(profile["fonts"]["heading"], "heading")
    title_font = resolve_pdf_font(profile["fonts"]["title"], "title")
    mono_font = resolve_pdf_font(profile["fonts"]["mono"], "mono")
    body_align = TA_CENTER if mode == "poetry" else TA_JUSTIFY if literary else TA_LEFT
    styles = {
        "title": ParagraphStyle("title", fontName=title_font, fontSize=30, leading=36,
                                 alignment=TA_CENTER, spaceAfter=18, textColor=theme["pdf_ink"]),
        "subtitle": ParagraphStyle("subtitle", fontName=body_font, fontSize=13, leading=18,
                                    alignment=TA_CENTER, textColor=theme["pdf_muted"], spaceAfter=6),
        "toc_entry": ParagraphStyle("toc_entry", fontName=body_font, fontSize=12, leading=20,
                                     alignment=TA_CENTER if literary else TA_LEFT,
                                     textColor=theme["pdf_ink"]),
        "chapter_head": ParagraphStyle("chapter_head", fontName=heading_font, fontSize=20,
                                        leading=24, alignment=TA_CENTER if literary else TA_LEFT,
                                        spaceBefore=10, spaceAfter=20, textColor=theme["pdf_ink"]),
        "chapter_num": ParagraphStyle("chapter_num", fontName=body_font, fontSize=10,
                                       leading=12, alignment=TA_CENTER if literary else TA_LEFT,
                                       textColor=theme["pdf_muted"], spaceAfter=4),
        "body": ParagraphStyle("body", fontName=body_font, fontSize=11.5, leading=18,
                                alignment=body_align,
                                firstLineIndent=(18 if literary and mode != "poetry" else 0),
                                spaceAfter=(4 if mode == "poetry" else 10),
                                textColor=theme["pdf_ink"]),
        "subheading": ParagraphStyle("subheading", fontName=heading_font, fontSize=13,
                                      leading=16, spaceBefore=10, spaceAfter=6, textColor=theme["pdf_accent"]),
        "bullet": ParagraphStyle("bullet", fontName=body_font, fontSize=11, leading=16,
                                  textColor=theme["pdf_ink"]),
        "table_cell": ParagraphStyle("table_cell", fontName=body_font, fontSize=9.5,
                                      leading=12, textColor=theme["pdf_ink"]),
        "table_header": ParagraphStyle("table_header", fontName=heading_font, fontSize=9.5,
                                        leading=12,
                                        textColor=theme.get("pdf_table_header_fg", rl_colors.white)),
        "mono": ParagraphStyle("mono", fontName=mono_font, fontSize=9.5, leading=12, textColor=theme["pdf_ink"]),
    }
    return styles


def build_pdf(project_id: str) -> str:
    proj = get_project(project_id)
    chapters = get_chapters(project_id)
    mode = proj["mode"]
    m = MODES.get(mode, MODES["novel"])
    doc_profile = project_document_profile(proj)
    theme = THEMES.get(proj["theme"], THEMES["classic_cream"])
    out_path = os.path.join(PDF_DIR, f"{project_id}.pdf")
    styles = build_pdf_styles(theme, mode, doc_profile)
    page_body_font = resolve_pdf_font(doc_profile["fonts"]["body"], "body")
    header_font = resolve_pdf_font(doc_profile["fonts"].get("caption", doc_profile["fonts"]["body"]), "body")
    title_font = resolve_pdf_font(doc_profile["fonts"]["title"], "title")
    heading_font = resolve_pdf_font(doc_profile["fonts"]["heading"], "heading")
    literary_doc = mode in ("novel", "short_story", "poetry") or doc_profile.get("document_type") == "narrative_nonfiction"
    show_section_numbers = bool(doc_profile.get("section_numbering")) or doc_profile.get("document_type") in {"legal_brief", "academic_paper"}

    def on_page(c: pdfcanvas.Canvas, doc):
        c.saveState()
        if doc.page > 1:
            c.setFont(header_font, 8)
            c.setFillColor(theme["pdf_muted"])
            c.drawString(doc.leftMargin, LETTER[1] - 0.45 * inch, proj["title"] or "Untitled")
            c.drawRightString(LETTER[0] - doc.rightMargin, LETTER[1] - 0.45 * inch, getattr(doc, "chapter_title", doc_profile["label"]))
            c.setStrokeColor(theme["pdf_muted"])
            c.setLineWidth(0.4)
            c.line(doc.leftMargin, LETTER[1] - 0.52 * inch, LETTER[0] - doc.rightMargin, LETTER[1] - 0.52 * inch)
            footer_y = 0.55 * inch
            if "confidential" in (proj.get("writing_notes") or "").lower() and not literary_doc:
                c.drawString(doc.leftMargin, footer_y, "CONFIDENTIAL")
            c.drawCentredString(LETTER[0] / 2, footer_y, str(doc.page - 1))
        if not literary_doc:
            c.setFillColor(theme["pdf_accent"])
            c.rect(0, LETTER[1] - 0.12 * inch, LETTER[0], 0.12 * inch, fill=1, stroke=0)
        c.restoreState()

    doc = BaseDocTemplate(out_path, pagesize=LETTER,
                           leftMargin=1.1 * inch, rightMargin=1.1 * inch,
                           topMargin=1 * inch, bottomMargin=1 * inch)
    frame = Frame(doc.leftMargin, doc.bottomMargin, doc.width, doc.height, id="normal")
    template = PageTemplate(id="main", frames=[frame], onPage=on_page)
    doc.addPageTemplates([template])

    story: List[Any] = []
    counters = {"figure": 0, "section": 0}

    # --- Title / cover page ---
    story.append(Spacer(1, 2.0 * inch))
    story.append(Paragraph(proj["title"] or "Untitled", styles["title"]))
    story.append(Paragraph(proj["genre"] or "", styles["subtitle"]))
    story.append(Spacer(1, 0.5 * inch))
    if m["literary"]:
        story.append(Paragraph(
            f"Generated {datetime.utcnow().strftime('%Y-%m-%d')} by an autonomous local "
            f"multi-agent writing system.", styles["subtitle"]
        ))
    else:
        story.append(Paragraph(
            f"Prepared {datetime.utcnow().strftime('%Y-%m-%d')} — {doc_profile['label']}", styles["subtitle"]
        ))
    story.append(PageBreak())

    if mode in ("novel", "short_story", "report") or len(chapters) > 1:
        story.append(Spacer(1, 0.4 * inch))
        story.append(Paragraph("Contents", ParagraphStyle("toc_title", parent=styles["chapter_head"], fontName=heading_font, fontSize=16, alignment=TA_CENTER, spaceAfter=12)))
        toc_rows = []
        for i, ch in enumerate(chapters):
            toc_rows.append([Paragraph(f"{i+1}. {escape_html(ch['title'] or m['unit'])}", styles["toc_entry"]), Paragraph(str(i + 2), ParagraphStyle("toc_page", parent=styles["toc_entry"], alignment=TA_RIGHT if 'TA_RIGHT' in globals() else TA_LEFT))])
        if toc_rows:
            toc_table = Table(toc_rows, colWidths=[5.6 * inch, 0.7 * inch])
            toc_table.setStyle(TableStyle([
                ("LINEBELOW", (0, 0), (-1, -1), 0.2, theme["pdf_muted"], [1, 3]),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 0),
                ("RIGHTPADDING", (0, 0), (-1, -1), 0),
            ]))
            story.append(toc_table)
        story.append(Paragraph("Page numbers are approximate in this build.", ParagraphStyle("toc_note", parent=styles["subtitle"], fontSize=8, alignment=TA_CENTER)))
        story.append(PageBreak())

    # --- Units ---
    for i, ch in enumerate(chapters):
        chapter_meta = ch.get("meta") or {}
        if isinstance(chapter_meta, str):
            try:
                chapter_meta = json.loads(chapter_meta)
            except Exception:
                chapter_meta = {}
        chapter_title = ch["title"] or f"{m['unit']} {i+1}"
        doc.chapter_title = chapter_title
        if literary_doc and doc_profile.get("document_type") in ("novel", "narrative_nonfiction", "short_story_collection"):
            opener_num = f"Chapter {i + 1}" if mode != "short_story" else f"Story {i + 1}"
            story.append(Spacer(1, 0.25 * inch))
            story.append(Paragraph(opener_num.upper() if mode != "short_story" else opener_num, ParagraphStyle("chapter_open_num", parent=styles["chapter_head"], fontName=title_font, fontSize=24, alignment=TA_CENTER, spaceAfter=8)))
            story.append(Paragraph(chapter_title, ParagraphStyle("chapter_open_title", parent=styles["chapter_head"], fontName=heading_font, fontSize=16, alignment=TA_CENTER, spaceAfter=8)))
            epigraph = _clean_text(chapter_meta.get("epigraph", ""))
            if epigraph:
                story.append(Paragraph(f"<i>{escape_html(epigraph)}</i>", ParagraphStyle("chapter_epigraph", parent=styles["body"], fontName=styles["body"].fontName, fontSize=10.5, leading=14, alignment=TA_CENTER, spaceAfter=12)))
            story.append(Spacer(1, 2))
            story.append(Paragraph("", ParagraphStyle("chapter_rule", parent=styles["body"], borderWidth=0.4, borderColor=theme["pdf_muted"], borderPadding=0, spaceAfter=10)))
        else:
            display_num = f"{i+1}." if show_section_numbers else f"{m['unit'].upper()} {i+1}"
            story.append(Paragraph(display_num, styles["chapter_num"]))
            story.append(Paragraph(chapter_title, styles["chapter_head"]))
        content = ch["content"] or ""

        if mode == "report":
            for btype, bval in parse_content_blocks(content):
                render_block_for_pdf(btype, bval, story, styles, theme, doc_profile, counters)
        elif mode == "poetry":
            for line in content.split("\n"):
                if line.strip():
                    story.append(Paragraph(escape_html(line), styles["body"]))
                else:
                    story.append(Spacer(1, 8))
        else:
            paragraphs = [p.strip() for p in re.split(r"\n\s*\n", content) if p.strip()]
            for idx_p, p in enumerate(paragraphs):
                story.append(Paragraph(escape_html(p), styles["body" if idx_p else "subheading"]))

        if i < len(chapters) - 1:
            if not m["literary"]:
                story.append(Spacer(1, 10))
                story.append(Paragraph("SECTION DIVIDER", styles["chapter_num"]))
            story.append(PageBreak())

    doc.build(story)
    logger.info("Built PDF for project %s at %s", project_id, out_path)
    return out_path


def escape_html(s: str) -> str:
    return (s or "").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


# =============================================================================
# FASTAPI APP
# =============================================================================

app = FastAPI(title="The Writer")
init_db()


class AgentSpec(BaseModel):
    name: str
    model: str
    role: str  # architect | writer | continuity

    @field_validator("role")
    @classmethod
    def validate_role(cls, v):
        if v not in ("architect", "writer", "continuity"):
            raise ValueError("role must be architect, writer, or continuity")
        return v


class CreateProjectRequest(BaseModel):
    title: str = Field(..., min_length=1)
    mode: str = Field("novel")
    genre: str = Field(..., min_length=1)
    premise: str = Field(..., min_length=1)
    document_type: str = ""
    num_units: int = Field(..., ge=MIN_UNITS, le=MAX_UNITS)
    words_per_unit: int = Field(..., ge=50, le=8000)
    theme: str = Field("classic_cream")
    writing_notes: str = ""
    embed_model: str = "nomic-embed-text"
    ocr_model: str = DEFAULT_OCR_MODEL
    agents: List[AgentSpec]

    @field_validator("mode")
    @classmethod
    def validate_mode(cls, v):
        if v not in MODES:
            raise ValueError(f"mode must be one of {list(MODES.keys())}")
        return v

    @field_validator("theme")
    @classmethod
    def validate_theme(cls, v):
        if v not in THEMES:
            raise ValueError(f"theme must be one of {list(THEMES.keys())}")
        return v

    @field_validator("document_type")
    @classmethod
    def validate_document_type(cls, v):
        if v and v not in DOCUMENT_TYPE_CHOICES:
            raise ValueError(f"document_type must be one of {DOCUMENT_TYPE_CHOICES} or empty")
        return v

    @field_validator("agents")
    @classmethod
    def validate_agents(cls, v):
        if not v:
            raise ValueError("At least one agent is required.")
        if not any(a.role == "architect" for a in v):
            raise ValueError("At least one architect agent is required.")
        if not any(a.role == "writer" for a in v):
            raise ValueError("At least one writer agent is required.")
        return v


@app.get("/api/health")
def api_health():
    ok, msg = ollama_health()
    return {"ollama_reachable": ok, "message": msg, "models": ollama_list_models()}


@app.get("/api/ollama/models")
def api_models():
    return {"models": ollama_list_models()}


@app.get("/api/modes")
def api_modes():
    return {"modes": MODES, "themes": {k: {"label": v["label"], "css_accent": v["css_accent"],
                                            "css_bg": v["css_bg"], "css_bg2": v["css_bg2"],
                                            "css_ink": v["css_ink"]} for k, v in THEMES.items()},
            "default_theme_by_mode": DEFAULT_THEME_BY_MODE,
            "document_types": DOCUMENT_TYPE_CHOICES,
            "report_document_types": REPORT_DOC_TYPES,
            "has_fitz": HAS_FITZ, "has_pypdf": HAS_PYPDF, "default_ocr_model": DEFAULT_OCR_MODEL}


@app.get("/api/projects")
def api_list_projects():
    return {"projects": list_projects()}


@app.post("/api/project")
def api_create_project(req: CreateProjectRequest):
    pid = str(uuid.uuid4())
    target_words = req.num_units * req.words_per_unit
    document_type = req.document_type.strip() or detect_document_type(
        req.mode, title=req.title, genre=req.genre, premise=req.premise, writing_notes=req.writing_notes
    )
    if req.mode != "report":
        document_type = MODE_TO_DEFAULT_DOC_TYPE.get(req.mode, document_type)
    with db() as conn:
        conn.execute(
            """INSERT INTO projects (id, title, mode, document_type, genre, premise, theme, writing_notes,
               num_units, words_per_unit, target_words, embed_model, ocr_model,
                    status, phase, created_at) VALUES (?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?)""",
            (pid, req.title, req.mode, document_type, req.genre, req.premise, req.theme, req.writing_notes,
             req.num_units, req.words_per_unit, target_words, req.embed_model, req.ocr_model,
             "draft", "idle", datetime.utcnow().isoformat())
        )
        for i, a in enumerate(req.agents):
            conn.execute(
                "INSERT INTO agents (id, project_id, name, model, role, order_idx) VALUES (?,?,?,?,?,?)",
                (str(uuid.uuid4()), pid, a.name, a.model, a.role, i)
            )
    PROGRESS[pid] = {"status": "draft", "phase": "idle", "total_units": 0,
                      "units_done": 0, "words_written": 0, "eta_seconds": 0}
    logger.info("Created project %s (%s, mode=%s, document_type=%s)", pid, req.title, req.mode, document_type)
    return {"project_id": pid}


@app.get("/api/project/{project_id}")
def api_get_project(project_id: str):
    proj = get_project(project_id)
    if not proj:
        return JSONResponse({"error": "not found"}, status_code=404)
    proj["agents"] = get_agents(project_id)
    return proj


@app.delete("/api/project/{project_id}")
def api_delete_project(project_id: str):
    STOP_FLAGS[project_id] = True
    delete_project(project_id)
    return {"status": "deleted"}


@app.post("/api/project/{project_id}/start")
def api_start(project_id: str):
    proj = get_project(project_id)
    if not proj:
        return JSONResponse({"error": "not found"}, status_code=404)

    if project_id in GEN_THREADS and GEN_THREADS[project_id].is_alive():
        return {"status": "already_running"}

    ok, msg = ollama_health()
    if not ok:
        return JSONResponse({"error": msg}, status_code=503)

    STOP_FLAGS[project_id] = False
    with db() as conn:
        conn.execute("UPDATE projects SET started_at=? WHERE id=?",
                     (datetime.utcnow().isoformat(), project_id))
    t = threading.Thread(target=run_generation, args=(project_id,), daemon=True)
    GEN_THREADS[project_id] = t
    t.start()
    logger.info("Started generation thread for project %s", project_id)
    return {"status": "started"}


@app.post("/api/project/{project_id}/stop")
def api_stop(project_id: str):
    STOP_FLAGS[project_id] = True
    log(project_id, "system", "system", "Stop requested by user.")
    return {"status": "stopping"}


@app.get("/api/project/{project_id}/status")
def api_status(project_id: str):
    proj = get_project(project_id)
    if not proj:
        return JSONResponse({"error": "not found"}, status_code=404)
    live = PROGRESS.get(project_id, {})
    merged = dict(proj)
    merged.update({k: v for k, v in live.items() if not k.startswith("_")})
    with db() as conn:
        logs = conn.execute(
            "SELECT ts, agent, role, message FROM logs WHERE project_id=? ORDER BY ts DESC LIMIT 50",
            (project_id,)
        ).fetchall()
    merged["recent_logs"] = [dict(r) for r in logs]
    merged["percent"] = round(
        100 * (merged.get("words_written", 0) / max(merged.get("target_words", 1), 1)), 1
    )
    merged["is_running"] = project_id in GEN_THREADS and GEN_THREADS[project_id].is_alive()
    return merged


@app.get("/api/project/{project_id}/chapters")
def api_chapters(project_id: str):
    with db() as conn:
        rows = conn.execute(
            "SELECT idx, title, synopsis, word_count, status, agent_name, revision_count "
            "FROM chapters WHERE project_id=? ORDER BY idx",
            (project_id,)
        ).fetchall()
    return {"chapters": [dict(r) for r in rows]}


@app.get("/api/project/{project_id}/chapter/{idx}")
def api_chapter(project_id: str, idx: int):
    row = get_chapter(project_id, idx)
    if not row:
        return JSONResponse({"error": "not found"}, status_code=404)
    return row


@app.get("/api/project/{project_id}/pdf")
def api_pdf(project_id: str):
    proj = get_project(project_id)
    if not proj:
        return JSONResponse({"error": "not found"}, status_code=404)
    chapters = get_chapters(project_id)
    if not any(c["content"] for c in chapters):
        return JSONResponse({"error": "No content written yet."}, status_code=400)
    path = build_pdf(project_id)
    safe_title = re.sub(r"[^a-zA-Z0-9_-]+", "_", proj["title"] or "document")
    return FileResponse(path, filename=f"{safe_title}.pdf", media_type="application/pdf")


# ---- Knowledge base endpoints (strictly scoped to one project each) --------

@app.get("/api/project/{project_id}/knowledge")
def api_kb_list(project_id: str):
    with db() as conn:
        rows = conn.execute(
            "SELECT id, filename, source_type, char_count, chunk_count, status, error, created_at "
            "FROM knowledge_docs WHERE project_id=? ORDER BY created_at DESC", (project_id,)
        ).fetchall()
    return {"docs": [dict(r) for r in rows]}


@app.post("/api/project/{project_id}/knowledge/upload")
async def api_kb_upload(project_id: str, file: UploadFile = File(...)):
    proj = get_project(project_id)
    if not proj:
        return JSONResponse({"error": "not found"}, status_code=404)
    filename = file.filename or "upload"
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
    raw = await file.read()

    if ext == "pdf":
        source_type = "pdf"
    elif ext in ("png", "jpg", "jpeg", "webp", "gif", "bmp"):
        source_type = "image"
    elif ext in ("txt", "md"):
        source_type = "text_file"
    else:
        return JSONResponse({"error": f"Unsupported file type '.{ext}'. Upload a PDF, image, "
                                       f".txt/.md file, or use the paste-text option."},
                             status_code=400)

    doc_id = str(uuid.uuid4())
    with db() as conn:
        conn.execute(
            "INSERT INTO knowledge_docs (id, project_id, filename, source_type, status, created_at) "
            "VALUES (?,?,?,?,?,?)",
            (doc_id, project_id, filename, source_type, "processing", datetime.utcnow().isoformat())
        )
    log(project_id, "system", "knowledge", f"Ingesting '{filename}' ({source_type})...")
    t = threading.Thread(
        target=ingest_knowledge_doc,
        args=(project_id, doc_id, filename, source_type, raw, None,
              proj["embed_model"], proj["ocr_model"] or DEFAULT_OCR_MODEL),
        daemon=True,
    )
    t.start()
    return {"doc_id": doc_id, "status": "processing"}


class PasteKnowledgeRequest(BaseModel):
    title: str = "Pasted text"
    text: str = Field(..., min_length=1)


class DetectDocTypeRequest(BaseModel):
    mode: str = Field("novel")
    title: str = ""
    genre: str = ""
    premise: str = ""
    writing_notes: str = ""

    @field_validator("mode")
    @classmethod
    def validate_mode(cls, v):
        if v not in MODES:
            raise ValueError(f"mode must be one of {list(MODES.keys())}")
        return v


@app.post("/api/project/{project_id}/knowledge/paste")
def api_kb_paste(project_id: str, req: PasteKnowledgeRequest):
    proj = get_project(project_id)
    if not proj:
        return JSONResponse({"error": "not found"}, status_code=404)
    doc_id = str(uuid.uuid4())
    with db() as conn:
        conn.execute(
            "INSERT INTO knowledge_docs (id, project_id, filename, source_type, status, created_at) "
            "VALUES (?,?,?,?,?,?)",
            (doc_id, project_id, req.title[:200], "paste", "processing", datetime.utcnow().isoformat())
        )
    t = threading.Thread(
        target=ingest_knowledge_doc,
        args=(project_id, doc_id, req.title[:200], "paste", None, req.text,
              proj["embed_model"], proj["ocr_model"] or DEFAULT_OCR_MODEL),
        daemon=True,
    )
    t.start()
    return {"doc_id": doc_id, "status": "processing"}


@app.delete("/api/project/{project_id}/knowledge/{doc_id}")
def api_kb_delete(project_id: str, doc_id: str):
    with db() as conn:
        conn.execute("DELETE FROM knowledge_docs WHERE id=? AND project_id=?", (doc_id, project_id))
        conn.execute("DELETE FROM vectors WHERE doc_id=? AND project_id=?", (doc_id, project_id))
    return {"status": "deleted"}


@app.post("/api/project/{project_id}/knowledge/upload-batch")
async def api_kb_upload_batch(project_id: str, files: List[UploadFile] = File(...)):
    proj = get_project(project_id)
    if not proj:
        return JSONResponse({"error": "not found"}, status_code=404)
    results = []
    for file in files:
        filename = file.filename or "upload"
        ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
        raw = await file.read()
        if ext == "pdf":
            source_type = "pdf"
        elif ext in ("png", "jpg", "jpeg", "webp", "gif", "bmp"):
            source_type = "image"
        elif ext in ("txt", "md"):
            source_type = "text_file"
        elif ext == "csv":
            source_type = "csv"
        elif ext in ("xlsx", "xlsm"):
            source_type = "xlsx"
        elif ext == "docx":
            source_type = "docx"
        elif ext == "pptx":
            source_type = "pptx"
        elif ext == "json":
            source_type = "json"
        else:
            results.append({"filename": filename, "status": "error", "error": f"Unsupported file type '.{ext}'"})
            continue
        doc_id = str(uuid.uuid4())
        with db() as conn:
            conn.execute(
                "INSERT INTO knowledge_docs (id, project_id, filename, source_type, status, created_at) VALUES (?,?,?,?,?,?)",
                (doc_id, project_id, filename, source_type, "processing", datetime.utcnow().isoformat())
            )
        log(project_id, "system", "knowledge", f"Ingesting '{filename}' ({source_type})...")
        t = threading.Thread(
            target=ingest_knowledge_doc,
            args=(project_id, doc_id, filename, source_type, raw, None,
                  proj["embed_model"], proj["ocr_model"] or DEFAULT_OCR_MODEL),
            daemon=True,
        )
        t.start()
        results.append({"doc_id": doc_id, "filename": filename, "status": "processing"})
    return {"results": results}


@app.get("/api/detect_doc_type")
def api_detect_doc_type(mode: str = "novel", title: str = "", genre: str = "", premise: str = "", writing_notes: str = ""):
    if mode not in MODES:
        return JSONResponse({"error": f"mode must be one of {list(MODES.keys())}"}, status_code=400)
    detected = detect_document_type(mode, title=title, genre=genre, premise=premise, writing_notes=writing_notes)
    profile = get_document_profile(mode, detected, title=title, genre=genre, premise=premise, writing_notes=writing_notes)
    return {"document_type": detected, "profile": profile}


@app.get("/api/project/{project_id}/knowledge/stats")
def api_kb_stats(project_id: str):
    with db() as conn:
        row = conn.execute(
            "SELECT COUNT(*) AS docs, COALESCE(SUM(char_count),0) AS chars, COALESCE(SUM(chunk_count),0) AS chunks FROM knowledge_docs WHERE project_id=?",
            (project_id,)
        ).fetchone()
        by_source = conn.execute(
            "SELECT source_type, COUNT(*) AS count FROM knowledge_docs WHERE project_id=? GROUP BY source_type ORDER BY count DESC",
            (project_id,)
        ).fetchall()
    return {"stats": dict(row), "by_source": [dict(r) for r in by_source]}


@app.delete("/api/project/{project_id}/knowledge")
def api_kb_clear(project_id: str):
    with db() as conn:
        doc_ids = [r[0] for r in conn.execute("SELECT id FROM knowledge_docs WHERE project_id=?", (project_id,)).fetchall()]
        conn.execute("DELETE FROM knowledge_docs WHERE project_id=?", (project_id,))
        conn.execute("DELETE FROM vectors WHERE project_id=? AND kind='knowledge'", (project_id,))
    return {"status": "cleared", "deleted_docs": len(doc_ids)}


@app.get("/api/project/{project_id}/knowledge/export", response_class=PlainTextResponse)
def api_kb_export(project_id: str):
    with db() as conn:
        rows = conn.execute(
            "SELECT filename, source_type, created_at FROM knowledge_docs WHERE project_id=? ORDER BY created_at DESC",
            (project_id,)
        ).fetchall()
        vectors = conn.execute(
            "SELECT text FROM vectors WHERE project_id=? AND kind='knowledge' ORDER BY chapter_idx ASC",
            (project_id,)
        ).fetchall()
    lines = [f"Knowledge Base Export for project {project_id}", ""]
    for row in rows:
        lines.append(f"- {row['filename']} [{row['source_type']}] {row['created_at']}")
    lines.append("")
    for idx, row in enumerate(vectors, 1):
        lines.append(f"[{idx}] {row['text']}")
        lines.append("")
    return "\n".join(lines)


# =============================================================================
# FRONTEND
# =============================================================================

HTML_PAGE = r"""
<!DOCTYPE html>
<html lang='en'>
<head>
<meta charset='UTF-8'>
<title>THE WRITER — Multi-Mode Autonomous Writing System</title>
<style>
  :root{
    --cream:#f1e9d8; --cream-2:#e8dec7; --panel:#faf5e9;
    --ink:#2a241c; --ink-soft:#5c5446;
    --rust:#a4502b; --rust-dark:#7e3c20; --line:#2a241c;
    --ok:#4a6b3a; --warn:#a4502b; --err:#9c2c2c;
  }
  *{box-sizing:border-box;}
  body{
    margin:0; background:var(--cream); color:var(--ink);
    font-family:'Courier New', monospace;
    background-image:repeating-linear-gradient(0deg, rgba(0,0,0,0.015) 0px, rgba(0,0,0,0.015) 1px, transparent 1px, transparent 3px);
  }
  header{ border-bottom:3px solid var(--line); padding:18px 28px; display:flex; justify-content:space-between; align-items:center; background:var(--cream-2); flex-wrap:wrap; gap:8px;}
  header h1{ margin:0; font-size:20px; letter-spacing:4px; font-weight:700; cursor:pointer;}
  header .tag{ font-size:11px; letter-spacing:2px; color:var(--ink-soft); }
  .wrap{ max-width:1320px; margin:0 auto; padding:24px; }
  .panel{ background:var(--panel); border:2px solid var(--line); box-shadow:6px 6px 0 var(--line); padding:20px; margin-bottom:22px; }
  .panel h2{ margin:0 0 16px 0; font-size:13px; letter-spacing:3px; text-transform:uppercase; border-bottom:2px solid var(--line); padding-bottom:8px; display:flex; justify-content:space-between; align-items:center;}
  label{ display:block; font-size:11px; letter-spacing:1px; text-transform:uppercase; color:var(--ink-soft); margin:14px 0 6px; }
  input[type=text], input[type=number], textarea, select{
    width:100%; background:#fffdf6; border:2px solid var(--line); color:var(--ink);
    font-family:'Courier New', monospace; font-size:14px; padding:9px 10px; border-radius:0; outline:none;
  }
  textarea{ resize:vertical; min-height:80px; }
  input:focus, textarea:focus, select:focus{ border-color:var(--rust); }
  .row{ display:flex; gap:14px; flex-wrap:wrap;}
  .row > div{ flex:1; min-width:180px;}
  button{
    font-family:'Courier New', monospace; font-weight:700; letter-spacing:2px; text-transform:uppercase; font-size:12px;
    background:var(--rust); color:#fff7ea; border:2px solid var(--line); box-shadow:4px 4px 0 var(--line);
    padding:11px 18px; cursor:pointer; border-radius:0; transition:transform .05s ease;
  }
  button:hover{ background:var(--rust-dark); }
  button:active{ transform:translate(4px,4px); box-shadow:0 0 0 var(--line); }
  button.ghost{ background:var(--panel); color:var(--ink); }
  button.danger{ background:var(--err); }
  button.small{ padding:6px 10px; font-size:10px; box-shadow:3px 3px 0 var(--line); }
  button:disabled{ opacity:0.4; cursor:not-allowed; }
  .mode-grid{ display:grid; grid-template-columns:repeat(auto-fit,minmax(200px,1fr)); gap:10px; margin-bottom:6px;}
  .mode-card{ border:2px solid var(--ink-soft); padding:12px; cursor:pointer; background:#fffdf6; }
  .mode-card.active{ border-color:var(--rust); background:#f1e3cf; box-shadow:3px 3px 0 var(--line); }
  .mode-card .t{ font-weight:700; font-size:13px; letter-spacing:1px; }
  .mode-card .d{ font-size:11px; color:var(--ink-soft); margin-top:4px; }
  .theme-grid{ display:flex; gap:10px; flex-wrap:wrap; }
  .theme-chip{ border:2px solid var(--ink-soft); padding:8px 12px; cursor:pointer; font-size:11px; display:flex; align-items:center; gap:8px; background:#fffdf6;}
  .theme-chip.active{ border-color:var(--rust); box-shadow:3px 3px 0 var(--line); }
  .theme-swatch{ width:14px; height:14px; display:inline-block; border:1px solid var(--line); }
  .agent-row{ display:flex; gap:10px; align-items:center; border:1px solid var(--ink-soft); padding:10px; margin-bottom:8px; background:#fffdf6; flex-wrap:wrap;}
  .agent-row .name{ width:140px; font-size:12px; font-weight:700; letter-spacing:1px; }
  .agent-row select{ flex:1; min-width:140px;}
  .remove-btn{ background:transparent; border:1px solid var(--err); color:var(--err); box-shadow:none; padding:6px 9px; }
  .remove-btn:active{ transform:none; }
  .screen{ display:none; }
  .screen.active{ display:block; }
  .grid{ display:grid; grid-template-columns: 280px 1fr 320px; gap:18px; align-items:start; }
  @media (max-width: 1000px){ .grid{ grid-template-columns: 1fr; } }
  .status-bar{ display:flex; gap:24px; flex-wrap:wrap; margin-bottom:10px; }
  .stat{ border:2px solid var(--line); padding:10px 14px; background:var(--panel); min-width:120px; }
  .stat .v{ font-size:22px; font-weight:700; }
  .stat .l{ font-size:10px; letter-spacing:1px; color:var(--ink-soft); text-transform:uppercase; }
  .progress-outer{ border:2px solid var(--line); height:22px; background:#fffdf6; margin:10px 0 18px; }
  .progress-inner{ height:100%; background:var(--rust); width:0%; transition:width .4s ease; }
  .chapter-item{ border:1px solid var(--ink-soft); padding:9px 10px; margin-bottom:6px; cursor:pointer; background:#fffdf6; font-size:12px; display:flex; justify-content:space-between; gap:8px;}
  .chapter-item:hover{ border-color:var(--rust); }
  .chapter-item.active{ border-color:var(--rust); background:#f1e3cf; }
  .badge{ font-size:9px; padding:2px 6px; border:1px solid var(--ink-soft); text-transform:uppercase; white-space:nowrap;}
  .badge.done{ background:#ddeacf; border-color:var(--ok); color:var(--ok); }
  .badge.writing{ background:#f6e2c4; border-color:var(--rust); color:var(--rust-dark); animation:pulse 1.4s infinite; }
  .badge.pending{ color:var(--ink-soft); }
  .badge.error{ background:#fbeaea; border-color:var(--err); color:var(--err); }
  .badge.ready{ background:#ddeacf; border-color:var(--ok); color:var(--ok); }
  .badge.processing{ background:#f6e2c4; border-color:var(--rust); color:var(--rust-dark); animation:pulse 1.4s infinite; }
  @keyframes pulse{ 0%,100%{opacity:1;} 50%{opacity:.4;} }
  #chapter-view{ white-space:pre-wrap; line-height:1.7; font-family:'Georgia','Times New Roman',serif; font-size:15px; max-height:70vh; overflow-y:auto; padding-right:6px;}
  #log-feed{ max-height:70vh; overflow-y:auto; font-size:11px; }
  .log-line{ border-bottom:1px dashed var(--ink-soft); padding:7px 0; }
  .log-line .who{ font-weight:700; color:var(--rust-dark); }
  .log-line .role{ color:var(--ink-soft); }
  .small{ font-size:11px; color:var(--ink-soft); }
  .err-box{ border:2px solid var(--err); color:var(--err); padding:10px; background:#fbeaea; margin-bottom:14px; font-size:12px; }
  .warn-box{ border:2px solid var(--warn); color:var(--rust-dark); padding:10px; background:#f6e2c4; margin-bottom:14px; font-size:12px; }
  a.download-link{ text-decoration:none; }
  table.project-table{ width:100%; border-collapse:collapse; font-size:12px; }
  table.project-table th, table.project-table td{ border:1px solid var(--ink-soft); padding:8px; text-align:left; }
  table.project-table th{ background:var(--cream-2); text-transform:uppercase; letter-spacing:1px; font-size:10px; }
  table.project-table tr:hover{ background:#fffdf6; }
  .kb-item{ display:flex; justify-content:space-between; align-items:center; gap:8px; border:1px solid var(--ink-soft); padding:8px 10px; margin-bottom:6px; background:#fffdf6; font-size:11px;}
  .tabs{ display:flex; gap:8px; margin-bottom:14px; }
  .tab-btn{ background:var(--panel); color:var(--ink); box-shadow:3px 3px 0 var(--line); }
  .tab-btn.active{ background:var(--rust); color:#fff7ea; }
  .filebtn{ display:inline-block; }
</style>
</head>
<body>

<header>
  <h1 onclick='showProjectList()'>THE WRITER</h1>
  <div class='tag'>LOCAL MULTI-MODE AUTONOMOUS WRITING SYSTEM — NO EXTERNAL APIs — ZERO HUMAN INPUT DURING WRITING</div>
</header>

<div class='wrap'>

  <div id='health-warning' class='warn-box' style='display:none;'></div>

  <!-- ============ PROJECT LIST SCREEN ============ -->
  <div id='list-screen' class='screen active'>
    <div class='panel'>
      <h2>My Books &amp; Documents <button onclick="showSetup()">+ NEW BOOK</button></h2>
      <table class='project-table' id='project-table'>
        <thead><tr><th>Title</th><th>Mode</th><th>Status</th><th>Progress</th><th>Units</th><th></th></tr></thead>
        <tbody id='project-table-body'></tbody>
      </table>
      <p class='small' id='no-projects-msg' style='display:none;'>Nothing yet. Click + NEW BOOK to start one.</p>
    </div>
  </div>

  <!-- ============ SETUP SCREEN ============ -->
  <div id='setup-screen' class='screen'>
    <div class='panel'>
      <h2>01 / Writing Mode <button class='ghost small' onclick='showProjectList()'>← BACK</button></h2>
      <div class='mode-grid' id='mode-grid'></div>
      <p class='small' id='mode-desc'></p>
    </div>

    <div class='panel'>
      <h2>02 / Configuration</h2>
      <div class='row'>
        <div><label id='label-title'>Title</label><input type='text' id='title' placeholder='THE LAST SIGNAL FROM KEPLER STATION'></div>
        <div><label id='label-genre'>Genre</label><input type='text' id='genre' value='Hard Science Fiction'></div>
      </div>
                        <div class='row'>
                                <div>
                                    <label>Document Type</label>
                                    <select id='document_type' onchange='updateDocTypeHint()'></select>
                                    <div class='small' id='doc-type-hint' style='margin-top:6px;color:var(--rust-dark);'>Auto-detect based on the title, genre, brief, and notes.</div>
                                </div>
                        </div>
      <label id='label-premise'>Premise / Brief (this is the ONLY human input — the rest is autonomous)</label>
      <textarea id='premise' placeholder='Describe the world, the core conflict, key characters, tone...'></textarea>
      <div class='row'>
        <div><label id='label-num-units'>Number of Chapters</label><input type='number' id='num_units' value='18' min='1' max='80'></div>
        <div><label id='label-words'>Target Words per Unit</label><input type='number' id='words_per_unit' value='2600' min='50' max='8000'></div>
        <div><label>Embedding Model (for local RAG memory)</label><select id='embed_model'></select></div>
      </div>
      <label>Author / House Style Notes (optional — applied to every agent prompt for this book)</label>
      <textarea id='writing_notes' placeholder='e.g. "Write in close third person, present tense" or "Use McKinsey-style pyramid structure, no hedging language"'></textarea>
    </div>

    <div class='panel'>
      <h2>03 / Theme &amp; Decor</h2>
      <p class='small'>Themes control typography and color throughout. Tables and report-style banners are only ever rendered in Executive Report mode — novels, short stories, and poetry stay in clean literary typesetting no matter which theme you pick.</p>
      <div class='theme-grid' id='theme-grid'></div>
    </div>

    <div class='panel'>
      <h2>04 / Knowledge Base <span class='small'>(strictly scoped to this book only)</span></h2>
      <p class='small'>Upload PDFs, images (scanned pages, screenshots, photos), or paste text directly. PDFs/images are OCR'd locally via <code id='ocr-model-label'></code>. This knowledge is only ever used for generating THIS book — it is never shared with any other project.</p>
      <p class='small' id='kb-hint' style='color:var(--rust-dark);'>Note: you can also add knowledge after creating the book, from its dashboard, before or during writing.</p>
    </div>

    <div class='panel'>
      <h2>05 / Agent Roster</h2>
      <p class='small' id='agent-desc'></p>
      <div id='agent-list'></div>
      <button class='ghost' onclick='addAgent()'>+ ADD AGENT</button>
    </div>

    <div class='panel'>
      <button onclick='createAndStart()' id='start-btn'>BEGIN AUTONOMOUS GENERATION</button>
      <span class='small' id='setup-error' style='margin-left:14px;color:var(--err);'></span>
    </div>
  </div>

  <!-- ============ DASHBOARD SCREEN ============ -->
  <div id='dash-screen' class='screen'>
    <div class='panel'>
      <div class='status-bar'>
        <div class='stat'><div class='v' id='stat-phase'>—</div><div class='l'>Phase</div></div>
        <div class='stat'><div class='v' id='stat-pct'>0%</div><div class='l'>Progress</div></div>
        <div class='stat'><div class='v' id='stat-chapters'>0/0</div><div class='l' id='stat-units-label'>Units</div></div>
        <div class='stat'><div class='v' id='stat-words'>0</div><div class='l'>Words Written</div></div>
        <div class='stat'><div class='v' id='stat-eta'>—</div><div class='l'>ETA</div></div>
        <div class='stat'>
          <button class='ghost small' onclick='showProjectList()' style='margin-bottom:6px;'>← MY BOOKS</button><br>
          <button class='danger small' onclick='stopGen()' id='stop-btn'>STOP</button>
          <button class='ghost small' onclick='resumeGen()' id='resume-btn' style='display:none;'>RESUME</button>
        </div>
        <div class='stat'>
          <a id='pdf-link' class='download-link' href='#' target='_blank'><button class='ghost'>DOWNLOAD PDF</button></a>
        </div>
      </div>
      <div class='progress-outer'><div class='progress-inner' id='progress-bar'></div></div>
      <div id='error-box' class='err-box' style='display:none;'></div>
    </div>

    <div class='tabs'>
      <button class='tab-btn active' onclick="showDashTab('write')" id='tab-write'>MANUSCRIPT</button>
      <button class='tab-btn' onclick="showDashTab('kb')" id='tab-kb'>KNOWLEDGE BASE</button>
    </div>

    <div id='dash-write'>
      <div class='grid'>
        <div class='panel'>
          <h2 id='panel-units-title'>Units</h2>
          <div id='chapter-list'></div>
        </div>
        <div class='panel'>
          <h2 id='chapter-title'>Select a unit</h2>
          <div id='chapter-view'></div>
        </div>
        <div class='panel'>
          <h2>Agent Activity</h2>
          <div id='log-feed'></div>
        </div>
      </div>
    </div>

    <div id='dash-kb' style='display:none;'>
      <div class='panel'>
        <h2>Knowledge Base for this Book</h2>
        <p class='small'>Scoped strictly to this book. Uploads are OCR'd locally (PDFs/images) via the project's OCR model, or ingested as-is (text files / pasted text).</p>
        <div class='row'>
          <div>
                        <label>Upload PDF / Image / Text File</label>
                        <input type='file' id='kb-file-input' accept='.pdf,.png,.jpg,.jpeg,.webp,.txt,.md,.csv,.xlsx,.docx,.pptx,.json' multiple>
                        <div class='row' style='margin-top:8px;gap:8px;'>
                            <button class='small' onclick='uploadKbFile()'>UPLOAD &amp; INGEST</button>
                            <button class='ghost small' onclick='clearKb()'>CLEAR KB</button>
                            <a id='kb-export-link' href='#' target='_blank' class='download-link'><button class='ghost small'>EXPORT TXT</button></a>
                        </div>
          </div>
          <div>
            <label>Or Paste Text Directly</label>
            <input type='text' id='kb-paste-title' placeholder='Title for this note'>
            <textarea id='kb-paste-text' placeholder='Paste reference text here...'></textarea>
            <button class='small' style='margin-top:8px;' onclick='pasteKb()'>ADD PASTED TEXT</button>
          </div>
        </div>
                <div class='small' id='kb-stats' style='margin-top:12px;color:var(--ink-soft);'></div>
        <div style='margin-top:16px;' id='kb-list'></div>
      </div>
    </div>
  </div>

</div>

<script>
let PROJECT_ID = null;
let SELECTED_CHAPTER = null;
let POLL_TIMER = null;
let KB_POLL_TIMER = null;
let agentCount = 0;
let availableModels = [];
let MODES_INFO = {};
let THEMES_INFO = {};
let DEFAULT_THEME_BY_MODE = {};
let DOCUMENT_TYPES = [];
let SELECTED_MODE = 'novel';
let SELECTED_THEME = 'classic_cream';
let DEFAULT_OCR_MODEL = 'glm-ocr:q8_0';
let DOC_TYPE_HINT_TIMER = null;

function show(id){
  document.querySelectorAll('.screen').forEach(s=>s.classList.remove('active'));
  document.getElementById(id).classList.add('active');
}

function showDashTab(which){
  document.getElementById('dash-write').style.display = which==='write' ? '' : 'none';
  document.getElementById('dash-kb').style.display = which==='kb' ? '' : 'none';
  document.getElementById('tab-write').classList.toggle('active', which==='write');
  document.getElementById('tab-kb').classList.toggle('active', which==='kb');
  if(which==='kb'){ refreshKb(); }
}

async function checkHealth(){
  try{
    const r = await fetch('/api/health');
    const d = await r.json();
    const box = document.getElementById('health-warning');
    if(!d.ollama_reachable){
      box.style.display = 'block';
      box.textContent = '⚠ ' + d.message + ' — generation will not start until Ollama is reachable.';
    } else if((d.models||[]).length === 0){
      box.style.display = 'block';
      box.textContent = '⚠ Ollama is reachable but no models are installed. Run: ollama pull llama3.1:8b';
    } else {
      box.style.display = 'none';
    }
    return d;
  }catch(e){
    document.getElementById('health-warning').style.display = 'block';
    document.getElementById('health-warning').textContent = '⚠ Could not reach the backend health check.';
    return {ollama_reachable:false, models:[]};
  }
}

async function loadModels(){
  const d = await checkHealth();
  availableModels = d.models || [];
  const embedSel = document.getElementById('embed_model');
  embedSel.innerHTML = '';
  const guesses = ['nomic-embed-text','mxbai-embed-large'];
  let opts = availableModels.length ? availableModels : guesses;
  opts.forEach(m=>{
    const o = document.createElement('option'); o.value=m; o.textContent=m;
    embedSel.appendChild(o);
  });
}

async function loadModesInfo(){
  const r = await fetch('/api/modes');
  const d = await r.json();
  MODES_INFO = d.modes; THEMES_INFO = d.themes; DEFAULT_THEME_BY_MODE = d.default_theme_by_mode;
    DOCUMENT_TYPES = d.document_types || [];
  DEFAULT_OCR_MODEL = d.default_ocr_model;
  document.getElementById('ocr-model-label').textContent = DEFAULT_OCR_MODEL;
  renderModeGrid();
  renderThemeGrid();
    renderDocumentTypeSelect();
}

function renderDocumentTypeSelect(){
    const select = document.getElementById('document_type');
    if(!select) return;
    const modeDocTypes = {
        novel: ['novel'],
        short_story: ['short_story_collection'],
        poetry: ['poetry_collection'],
        report: DOCUMENT_TYPES.filter(t => t !== 'novel' && t !== 'short_story_collection' && t !== 'poetry_collection')
    };
    const options = ['auto'];
    (modeDocTypes[SELECTED_MODE] || DOCUMENT_TYPES).forEach(t => {
        if(!options.includes(t)) options.push(t);
    });
    select.innerHTML = options.map(t => {
        const label = t === 'auto' ? 'Auto-detect from title / brief' : t.replace(/_/g, ' ').replace(/\b\w/g, c => c.toUpperCase());
        return `<option value="${t}" ${t === (select.value || 'auto') ? 'selected' : ''}>${label}</option>`;
    }).join('');
    if(!select.value){ select.value = 'auto'; }
}

async function updateDocTypeHint(){
    const hint = document.getElementById('doc-type-hint');
    if(!hint) return;
    if(DOC_TYPE_HINT_TIMER){ clearTimeout(DOC_TYPE_HINT_TIMER); }
    DOC_TYPE_HINT_TIMER = setTimeout(async ()=>{
        try{
            const params = new URLSearchParams({
                mode: SELECTED_MODE,
                title: document.getElementById('title').value || '',
                genre: document.getElementById('genre').value || '',
                premise: document.getElementById('premise').value || '',
                writing_notes: document.getElementById('writing_notes').value || ''
            });
            const r = await fetch('/api/detect_doc_type?' + params.toString());
            const d = await r.json();
            if(d.error){
                hint.textContent = 'Auto-detect unavailable.';
                return;
            }
            const profile = d.profile || {};
            hint.textContent = 'Detected: ' + (profile.label || d.document_type || 'auto') + ' — ' + (profile.description || 'profile applied automatically');
            const select = document.getElementById('document_type');
            if(select && (!select.value || select.value === 'auto')){
                const options = Array.from(select.options || []);
                const hit = options.find(o => o.value === d.document_type);
                if(hit){ select.value = d.document_type; }
            }
        }catch(e){
            hint.textContent = 'Auto-detect unavailable.';
        }
    }, 250);
}

function renderModeGrid(){
  const grid = document.getElementById('mode-grid');
  grid.innerHTML = '';
  const descs = {
    novel: 'Multi-chapter long-form fiction with full continuity and a story bible.',
    short_story: 'A collection of standalone short stories sharing a premise.',
    poetry: 'A collection of poems, form-aware, one poem per unit.',
    report: 'A professional executive report with formal sections and data tables.'
  };
  Object.keys(MODES_INFO).forEach(key=>{
    const info = MODES_INFO[key];
    const card = document.createElement('div');
    card.className = 'mode-card' + (key===SELECTED_MODE ? ' active' : '');
    card.innerHTML = `<div class='t'>${info.label}</div><div class='d'>${descs[key]||''}</div>`;
    card.onclick = ()=> selectMode(key);
    grid.appendChild(card);
  });
  document.getElementById('mode-desc').textContent = descs[SELECTED_MODE] || '';
}

function renderThemeGrid(){
  const grid = document.getElementById('theme-grid');
  grid.innerHTML = '';
  Object.keys(THEMES_INFO).forEach(key=>{
    const info = THEMES_INFO[key];
    const chip = document.createElement('div');
    chip.className = 'theme-chip' + (key===SELECTED_THEME ? ' active' : '');
    chip.innerHTML = `<span class='theme-swatch' style='background:${info.css_accent}'></span>${info.label}`;
    chip.onclick = ()=> { SELECTED_THEME = key; renderThemeGrid(); };
    grid.appendChild(chip);
  });
}

function selectMode(key){
  SELECTED_MODE = key;
  SELECTED_THEME = DEFAULT_THEME_BY_MODE[key] || SELECTED_THEME;
  const info = MODES_INFO[key];
  document.getElementById('label-num-units').textContent = 'Number of ' + info.unit_plural;
  document.getElementById('num_units').value = info.default_units;
  document.getElementById('words_per_unit').value = info.default_words_per_unit;
  if(key === 'report'){
    document.getElementById('label-title').textContent = 'Report Title';
    document.getElementById('label-genre').textContent = 'Domain / Industry';
    document.getElementById('label-premise').textContent = 'Purpose / Brief (the only human input)';
    document.getElementById('genre').value = 'Enterprise Technology';
    document.getElementById('agent-desc').textContent = "ARCHITECT designs the section structure (and marks which sections need data tables). WRITER agents rotate sections. REVIEWER agents (role 'continuity') check factual/tonal consistency and trigger one corrective pass when real issues are found.";
  } else if(key === 'poetry'){
    document.getElementById('label-title').textContent = 'Collection Title';
    document.getElementById('label-genre').textContent = 'Tradition / Style';
    document.getElementById('label-premise').textContent = 'Premise / Unifying Theme (the only human input)';
    document.getElementById('genre').value = 'Contemporary Free Verse';
    document.getElementById('agent-desc').textContent = "ARCHITECT designs the poem-by-poem table of contents (with form per poem). WRITER agents rotate poems. REVIEWER agents optionally check form/subject drift.";
  } else if(key === 'short_story'){
    document.getElementById('label-title').textContent = 'Collection Title';
    document.getElementById('label-genre').textContent = 'Genre';
    document.getElementById('label-premise').textContent = 'Premise / Unifying Idea (the only human input)';
    document.getElementById('genre').value = 'Literary Fiction';
    document.getElementById('agent-desc').textContent = "ARCHITECT designs the story-by-story table of contents. WRITER agents rotate stories. REVIEWER agents check internal consistency and trigger one corrective pass when real issues are found.";
  } else {
    document.getElementById('label-title').textContent = 'Title';
    document.getElementById('label-genre').textContent = 'Genre';
    document.getElementById('label-premise').textContent = 'Premise / Seed Context (the only human input)';
    document.getElementById('genre').value = 'Hard Science Fiction';
    document.getElementById('agent-desc').textContent = "ARCHITECT designs the chapter outline (first one is used). WRITER agents rotate chapters. REVIEWER agents (role 'continuity') review for consistency and trigger one corrective pass when real issues are found.";
  }
  renderModeGrid();
  renderThemeGrid();
    renderDocumentTypeSelect();
    updateDocTypeHint();
}

function modelOptionsHTML(selected){
  let opts = availableModels.length ? availableModels : ['llama3.1:8b'];
  return opts.map(m=>`<option value="${m}" ${m===selected?'selected':''}>${m}</option>`).join('');
}

function addAgent(role, modelGuess){
  agentCount++;
  const id = 'agent_' + agentCount;
  role = role || 'writer';
  const div = document.createElement('div');
  div.className = 'agent-row';
  div.id = id;
  div.innerHTML = `
    <div class='name'>${role.toUpperCase()}-${agentCount}</div>
    <select class='model-select'>${modelOptionsHTML(modelGuess)}</select>
    <select class='role-select'>
      <option value='architect' ${role==='architect'?'selected':''}>Architect</option>
      <option value='writer' ${role==='writer'?'selected':''}>Writer</option>
      <option value='continuity' ${role==='continuity'?'selected':''}>Reviewer</option>
    </select>
    <button class='remove-btn' onclick="document.getElementById('${id}').remove()">✕</button>
  `;
  document.getElementById('agent-list').appendChild(div);
}

async function showSetup(){
  document.getElementById('agent-list').innerHTML = '';
  agentCount = 0;
  await loadModesInfo();
  selectMode(SELECTED_MODE);
  await loadModels();
  addAgent('architect');
  addAgent('writer');
  addAgent('writer');
  addAgent('continuity');
    ['title','genre','premise','writing_notes'].forEach(id=>{
        const el = document.getElementById(id);
        if(el){ el.oninput = updateDocTypeHint; }
    });
  show('setup-screen');
}

async function createAndStart(){
  const title = document.getElementById('title').value.trim() || 'Untitled';
  const genre = document.getElementById('genre').value.trim() || 'General';
  const premise = document.getElementById('premise').value.trim();
  const num_units = parseInt(document.getElementById('num_units').value) || MODES_INFO[SELECTED_MODE].default_units;
  const words_per_unit = parseInt(document.getElementById('words_per_unit').value) || MODES_INFO[SELECTED_MODE].default_words_per_unit;
  const embed_model = document.getElementById('embed_model').value;
  const writing_notes = document.getElementById('writing_notes').value.trim();
    const document_type = (document.getElementById('document_type').value || 'auto') === 'auto' ? '' : document.getElementById('document_type').value;
  const errEl = document.getElementById('setup-error');
  errEl.textContent = '';

  if(!premise){ errEl.textContent = 'Premise / brief is required.'; return; }

  const rows = document.querySelectorAll('.agent-row');
  if(rows.length === 0){ errEl.textContent = 'Add at least one architect and one writer agent.'; return; }
  const agents = [];
  rows.forEach(r=>{
    agents.push({
      name: r.querySelector('.name').textContent,
      model: r.querySelector('.model-select').value,
      role: r.querySelector('.role-select').value
    });
  });
  if(!agents.some(a=>a.role==='architect')){ errEl.textContent = 'Need at least one Architect agent.'; return; }
  if(!agents.some(a=>a.role==='writer')){ errEl.textContent = 'Need at least one Writer agent.'; return; }

  document.getElementById('start-btn').disabled = true;
  try{
    const r = await fetch('/api/project', {
      method:'POST', headers:{'Content-Type':'application/json'},
      body: JSON.stringify({
                title, mode: SELECTED_MODE, genre, premise, document_type, num_units, words_per_unit,
        theme: SELECTED_THEME, writing_notes, embed_model,
        ocr_model: DEFAULT_OCR_MODEL, agents
      })
    });
    const d = await r.json();
    if(d.error || d.detail){
      errEl.textContent = d.error || JSON.stringify(d.detail);
      document.getElementById('start-btn').disabled=false; return;
    }
    PROJECT_ID = d.project_id;
    const startR = await fetch(`/api/project/${PROJECT_ID}/start`, {method:'POST'});
    const startD = await startR.json();
    if(startD.error){
      errEl.textContent = startD.error;
      document.getElementById('start-btn').disabled=false; return;
    }
    openDashboard(PROJECT_ID);
  }catch(e){
    errEl.textContent = 'Failed to start: ' + e;
    document.getElementById('start-btn').disabled = false;
  }
}

async function openDashboard(pid){
  PROJECT_ID = pid;
  SELECTED_CHAPTER = null;
  document.getElementById('chapter-list').innerHTML = '';
  document.getElementById('chapter-view').textContent = '';
  document.getElementById('chapter-title').textContent = 'Select a unit';
  document.getElementById('pdf-link').href = `/api/project/${pid}/pdf`;
  showDashTab('write');
  show('dash-screen');
  const r = await fetch(`/api/project/${pid}`);
  const proj = await r.json();
  if(!MODES_INFO || Object.keys(MODES_INFO).length===0){ await loadModesInfo(); }
  const info = MODES_INFO[proj.mode] || MODES_INFO['novel'];
  document.getElementById('panel-units-title').textContent = info.unit_plural;
  document.getElementById('stat-units-label').textContent = info.unit_plural;
  startPolling();
}

async function stopGen(){
  if(!PROJECT_ID) return;
  await fetch(`/api/project/${PROJECT_ID}/stop`, {method:'POST'});
}

async function resumeGen(){
  if(!PROJECT_ID) return;
  document.getElementById('resume-btn').disabled = true;
  await fetch(`/api/project/${PROJECT_ID}/start`, {method:'POST'});
}

function fmtETA(sec){
  if(!sec || sec<=0) return '—';
  const h = Math.floor(sec/3600), m = Math.floor((sec%3600)/60);
  if(h>0) return `${h}h ${m}m`;
  return `${m}m`;
}

async function poll(){
  if(!PROJECT_ID) return;
  try{
    const r = await fetch(`/api/project/${PROJECT_ID}/status`);
    const d = await r.json();
    if(d.error){ return; }
    document.getElementById('stat-phase').textContent = (d.phase||'—').toUpperCase();
    document.getElementById('stat-pct').textContent = (d.percent||0) + '%';
    document.getElementById('stat-chapters').textContent = `${d.units_done||0}/${d.total_units||0}`;
    document.getElementById('stat-words').textContent = (d.words_written||0).toLocaleString();
    document.getElementById('stat-eta').textContent = fmtETA(d.eta_seconds);
    document.getElementById('progress-bar').style.width = Math.min(d.percent||0,100) + '%';

    const stopBtn = document.getElementById('stop-btn');
    const resumeBtn = document.getElementById('resume-btn');
    if(d.status === 'stopped' || d.status === 'error'){
      stopBtn.style.display = 'none';
      resumeBtn.style.display = 'inline-block';
      resumeBtn.disabled = false;
    } else {
      stopBtn.style.display = 'inline-block';
      resumeBtn.style.display = 'none';
    }

    const errBox = document.getElementById('error-box');
    if(d.status === 'error'){
      errBox.style.display = 'block';
      errBox.textContent = 'ERROR: ' + (d.error || 'unknown error') + ' — fix the issue (e.g. start Ollama, pull the model) and click RESUME to continue from where it left off.';
    } else {
      errBox.style.display = 'none';
    }

    const feed = document.getElementById('log-feed');
    feed.innerHTML = (d.recent_logs||[]).map(l=>
      `<div class='log-line'><span class='who'>${escapeHtml(l.agent)}</span> <span class='role'>[${escapeHtml(l.role)}]</span><br>${escapeHtml(l.message)}</div>`
    ).join('');

    await refreshChapters();
  }catch(e){ /* keep polling */ }
}

async function refreshChapters(){
  const r = await fetch(`/api/project/${PROJECT_ID}/chapters`);
  const d = await r.json();
  const list = document.getElementById('chapter-list');
  list.innerHTML = '';
  (d.chapters||[]).forEach(ch=>{
    const item = document.createElement('div');
    item.className = 'chapter-item' + (SELECTED_CHAPTER===ch.idx ? ' active':'');
    item.innerHTML = `<span>${ch.idx+1}. ${escapeHtml(ch.title||'')}</span><span class='badge ${ch.status}'>${ch.status}${ch.revision_count?' ('+ch.revision_count+' fix)':''}</span>`;
    item.onclick = ()=> selectChapter(ch.idx, ch.title);
    list.appendChild(item);
  });
  if(SELECTED_CHAPTER !== null){
    await loadChapterText(SELECTED_CHAPTER);
  }
}

async function selectChapter(idx, title){
  SELECTED_CHAPTER = idx;
  document.getElementById('chapter-title').textContent = `${idx+1}. ${title}`;
  await loadChapterText(idx);
  await refreshChapters();
}

async function loadChapterText(idx){
  const r = await fetch(`/api/project/${PROJECT_ID}/chapter/${idx}`);
  const d = await r.json();
  document.getElementById('chapter-view').textContent = d.content || '(writing not yet started for this unit)';
}

function escapeHtml(s){
  return (s||'').toString().replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;');
}

function startPolling(){
  if(POLL_TIMER) clearInterval(POLL_TIMER);
  poll();
  POLL_TIMER = setInterval(poll, 3000);
}

async function refreshKb(){
  if(!PROJECT_ID) return;
    const [docsR, statsR] = await Promise.all([
        fetch(`/api/project/${PROJECT_ID}/knowledge`),
        fetch(`/api/project/${PROJECT_ID}/knowledge/stats`)
    ]);
    const d = await docsR.json();
    const s = await statsR.json();
  const list = document.getElementById('kb-list');
  list.innerHTML = '';
  (d.docs||[]).forEach(doc=>{
    const item = document.createElement('div');
    item.className = 'kb-item';
    item.innerHTML = `<span>${escapeHtml(doc.filename)} <span class='small'>(${escapeHtml(doc.source_type)}, ${doc.chunk_count||0} chunks)</span></span>
      <span><span class='badge ${doc.status}'>${doc.status}</span> <button class='remove-btn' onclick="deleteKb('${doc.id}')">✕</button></span>`;
    if(doc.status==='error'){ item.title = doc.error; }
    list.appendChild(item);
  });
    const stats = s.stats || {};
    document.getElementById('kb-stats').textContent = `Docs: ${stats.docs||0} · Chars: ${(stats.chars||0).toLocaleString()} · Chunks: ${stats.chunks||0}`;
    const exportLink = document.getElementById('kb-export-link');
    if(exportLink){ exportLink.href = `/api/project/${PROJECT_ID}/knowledge/export`; }
}

async function uploadKbFile(){
  const input = document.getElementById('kb-file-input');
  if(!input.files.length){ return; }
  const fd = new FormData();
    Array.from(input.files).forEach(file=> fd.append('files', file));
    const endpoint = input.files.length > 1 ? `/api/project/${PROJECT_ID}/knowledge/upload-batch` : `/api/project/${PROJECT_ID}/knowledge/upload-batch`;
    await fetch(endpoint, {method:'POST', body: fd});
  input.value = '';
  await refreshKb();
  if(!KB_POLL_TIMER){ KB_POLL_TIMER = setInterval(refreshKb, 4000); }
}

async function pasteKb(){
  const title = document.getElementById('kb-paste-title').value.trim() || 'Pasted text';
  const text = document.getElementById('kb-paste-text').value.trim();
  if(!text) return;
  await fetch(`/api/project/${PROJECT_ID}/knowledge/paste`, {
    method:'POST', headers:{'Content-Type':'application/json'},
    body: JSON.stringify({title, text})
  });
  document.getElementById('kb-paste-title').value = '';
  document.getElementById('kb-paste-text').value = '';
  await refreshKb();
}

async function deleteKb(docId){
  await fetch(`/api/project/${PROJECT_ID}/knowledge/${docId}`, {method:'DELETE'});
  await refreshKb();
}

async function clearKb(){
    if(!confirm('Clear all knowledge documents for this book?')) return;
    await fetch(`/api/project/${PROJECT_ID}/knowledge`, {method:'DELETE'});
    await refreshKb();
}

async function showProjectList(){
  if(POLL_TIMER){ clearInterval(POLL_TIMER); POLL_TIMER = null; }
  if(KB_POLL_TIMER){ clearInterval(KB_POLL_TIMER); KB_POLL_TIMER = null; }
  await checkHealth();
  if(!MODES_INFO || Object.keys(MODES_INFO).length===0){ await loadModesInfo(); }
  const r = await fetch('/api/projects');
  const d = await r.json();
  const tbody = document.getElementById('project-table-body');
  tbody.innerHTML = '';
  document.getElementById('no-projects-msg').style.display = (d.projects||[]).length ? 'none' : 'block';
  (d.projects||[]).forEach(p=>{
    const tr = document.createElement('tr');
    const pct = p.target_words ? Math.round(100*(p.words_written||0)/p.target_words) : 0;
    const modeLabel = (MODES_INFO[p.mode]||{}).label || p.mode;
    tr.innerHTML = `
      <td>${escapeHtml(p.title)}</td>
      <td>${escapeHtml(modeLabel)}</td>
      <td>${escapeHtml(p.status)} / ${escapeHtml(p.phase)}</td>
      <td>${p.units_done||0}/${p.total_units||0} · ${pct}%</td>
      <td>${p.num_units||0}</td>
      <td>
        <button class='small' onclick="openDashboard('${p.id}')">OPEN</button>
        <a href='/api/project/${p.id}/pdf' target='_blank'><button class='ghost small'>PDF</button></a>
        <button class='danger small' onclick="deleteProject('${p.id}')">DEL</button>
      </td>
    `;
    tbody.appendChild(tr);
  });
  show('list-screen');
}

async function deleteProject(pid){
  if(!confirm('Delete this project permanently? This cannot be undone.')) return;
  await fetch(`/api/project/${pid}`, {method:'DELETE'});
  await showProjectList();
}

showProjectList();
</script>

</body>
</html>
"""


@app.get("/", response_class=HTMLResponse)
def index():
    return HTML_PAGE


# =============================================================================
# ENTRY POINT
# =============================================================================

if __name__ == "__main__":
    import multiprocessing
    multiprocessing.freeze_support()  # no-op unless frozen; safe/required for PyInstaller on Windows

    logger.info("DB at: %s", DB_PATH)
    logger.info("PDFs at: %s", PDF_DIR)
    logger.info("Logs at: %s", LOG_PATH)
    logger.info("Ollama endpoint: %s", OLLAMA_URL)
    logger.info("PyMuPDF available: %s | pypdf available: %s", HAS_FITZ, HAS_PYPDF)
    ok, msg = ollama_health()
    if ok:
        logger.info("Ollama health check: OK (%s)", msg)
    else:
        logger.warning("Ollama health check: FAILED (%s) — fix this before starting a book.", msg)

    # When running as a packaged .exe, auto-open the UI in the default browser
    # since there's no terminal the user is expected to read a URL from.
    if getattr(sys, "frozen", False):
        import webbrowser
        threading.Timer(1.2, lambda: webbrowser.open(f"http://127.0.0.1:{PORT}")).start()

    uvicorn.run(app, host="0.0.0.0", port=PORT)
